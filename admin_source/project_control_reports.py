# -*- coding: utf-8 -*-
"""گزارش‌های برنامه عملیاتی و کنترل پروژه نسخه ۶.۷."""

from jalali_utils import convert_dates_in_text, format_jalali, now_jalali, jalali_year, install_openpyxl_jalali_patch, install_pptx_jalali_patch
install_openpyxl_jalali_patch()
install_pptx_jalali_patch()
from datetime import datetime, timedelta


def _safe(value):
    return "—" if value in (None, "") else value


def export_project_control_excel(db, path, fiscal_year=None, zone_id=None):
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.utils import get_column_letter

    summary = db.get_project_control_summary(fiscal_year=fiscal_year, zone_id=zone_id)
    programs = db.get_annual_programs(fiscal_year=fiscal_year, zone_id=zone_id)
    projects = db.get_projects(zone_id=zone_id)
    if fiscal_year:
        projects = [p for p in projects if str(p.get("fiscal_year") or "") == str(fiscal_year)]
    project_ids = {p["id"] for p in projects}
    milestones = [m for m in db.get_project_milestones() if m.get("project_id") in project_ids]
    indicators = [i for i in db.get_project_indicators() if i.get("project_id") in project_ids or i.get("program_id") in {p["id"] for p in programs}]
    risks = [r for r in db.get_project_risks(zone_id=zone_id) if not project_ids or r.get("project_id") in project_ids or r.get("project_id") is None]
    changes = [c for c in db.get_project_change_requests() if zone_id is None or c.get("zone_id") == zone_id]
    alerts = db.get_project_control_alerts(fiscal_year=fiscal_year, zone_id=zone_id)

    wb = Workbook()
    ws = wb.active
    ws.title = "داشبورد"
    navy = "13294B"; gold = "C9A227"; white = "FFFFFF"; light = "EEF2F7"; red = "F8D7DA"
    thin = Side(style="thin", color="D5DCE6")

    def style_sheet(sheet, widths=None):
        sheet.sheet_view.rightToLeft = True
        sheet.freeze_panes = "A2"
        for cell in sheet[1]:
            cell.font = Font(bold=True, color=white)
            cell.fill = PatternFill("solid", fgColor=navy)
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell.border = Border(bottom=thin)
        for row in sheet.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(horizontal="right", vertical="center", wrap_text=True)
                cell.border = Border(bottom=thin)
        if widths:
            for idx, width in enumerate(widths, 1):
                sheet.column_dimensions[get_column_letter(idx)].width = width

    ws.merge_cells("A1:D1")
    ws["A1"] = "داشبورد برنامه عملیاتی و کنترل پروژه"
    ws["A1"].font = Font(bold=True, color=white, size=15)
    ws["A1"].fill = PatternFill("solid", fgColor=navy)
    ws["A1"].alignment = Alignment(horizontal="center")
    ws.append(["سال مالی", fiscal_year or "همه", "تاریخ تهیه", format_jalali(datetime.now().strftime("%Y-%m-%d"))])
    ws.append([])
    metrics = [
        ("تعداد برنامه‌ها", summary["programs_count"]),
        ("تعداد پروژه‌ها", summary["projects_count"]),
        ("پروژه فعال", summary["active_projects"]),
        ("پروژه معوق", summary["overdue_projects"]),
        ("بودجه برنامه‌ریزی‌شده", summary["planned_budget"]),
        ("هزینه واقعی", summary["actual_cost"]),
        ("انحراف هزینه", summary["cost_variance"]),
        ("میانگین پیشرفت", summary["average_progress"]),
        ("تحقق شاخص‌ها", summary["indicator_achievement"]),
        ("ریسک بالا", summary["high_risks"]),
        ("تغییر در انتظار", summary["pending_changes"]),
        ("هشدار باز", summary["alerts_count"]),
    ]
    ws.append(["شاخص", "مقدار", "شاخص", "مقدار"])
    for i in range(0, len(metrics), 2):
        left = metrics[i]
        right = metrics[i + 1] if i + 1 < len(metrics) else ("", "")
        ws.append([left[0], left[1], right[0], right[1]])
    for cell in ws[4]:
        cell.font = Font(bold=True, color=white); cell.fill = PatternFill("solid", fgColor=gold)
    for row in ws.iter_rows(min_row=2):
        for cell in row: cell.alignment = Alignment(horizontal="right", vertical="center")
    for col, width in zip("ABCD", [32, 18, 32, 18]): ws.column_dimensions[col].width = width

    sh = wb.create_sheet("برنامه‌ها")
    sh.append(["سال مالی", "عنوان", "هدف راهبردی", "بلوک", "دستگاه مسئول", "مدیر برنامه", "شروع", "پایان", "بودجه مصوب", "پیشرفت٪", "وضعیت", "تعداد پروژه"])
    for p in programs:
        sh.append([p.get("fiscal_year"), p.get("title"), p.get("strategic_goal"), p.get("zone_name"), p.get("responsible_agency"), p.get("program_manager"), p.get("start_date"), p.get("end_date"), p.get("approved_budget"), p.get("progress_percent"), p.get("status"), p.get("project_count")])
    style_sheet(sh, [13, 36, 40, 22, 24, 22, 14, 14, 18, 14, 18, 14])

    sh = wb.create_sheet("پروژه‌ها")
    sh.append(["کد", "عنوان", "برنامه", "بلوک", "دستگاه", "مدیر پروژه", "شروع", "پایان", "بودجه", "هزینه", "برنامه٪", "واقعی٪", "انحراف پیشرفت", "انحراف هزینه", "اولویت", "وضعیت"])
    for p in projects:
        sh.append([p.get("project_code"), p.get("title"), p.get("program_title"), p.get("zone_name"), p.get("responsible_agency"), p.get("project_manager"), p.get("start_date"), p.get("end_date"), p.get("planned_budget"), p.get("actual_cost"), p.get("planned_progress"), p.get("actual_progress"), p.get("progress_variance"), p.get("cost_variance"), p.get("priority"), p.get("status")])
    style_sheet(sh, [16, 36, 32, 20, 24, 22, 13, 13, 16, 16, 12, 12, 16, 16, 12, 18])
    for row in range(2, sh.max_row + 1):
        if sh.cell(row, 14).value and sh.cell(row, 14).value > 0:
            sh.cell(row, 14).fill = PatternFill("solid", fgColor=red)

    sh = wb.create_sheet("نقاط عطف")
    sh.append(["کد پروژه", "پروژه", "عنوان نقطه عطف", "بلوک", "سررسید", "تکمیل", "وزن", "وضعیت", "معوق"])
    for m in milestones:
        sh.append([m.get("project_code"), m.get("project_title"), m.get("title"), m.get("zone_name"), m.get("due_date"), m.get("completed_date"), m.get("weight"), m.get("status"), "بله" if m.get("is_overdue") else "خیر"])
    style_sheet(sh, [16, 32, 36, 20, 14, 14, 10, 18, 10])

    sh = wb.create_sheet("شاخص‌ها")
    sh.append(["برنامه", "پروژه", "عنوان شاخص", "واحد", "مبنا", "هدف", "عملکرد", "جهت", "تحقق٪", "تاریخ اندازه‌گیری", "یادداشت"])
    for i in indicators:
        sh.append([i.get("program_title"), i.get("project_title"), i.get("title"), i.get("unit"), i.get("baseline_value"), i.get("target_value"), i.get("actual_value"), i.get("direction"), i.get("achievement_percent"), i.get("measurement_date"), i.get("notes")])
    style_sheet(sh, [30, 30, 36, 12, 12, 12, 12, 12, 12, 16, 40])

    sh = wb.create_sheet("ریسک‌ها")
    sh.append(["برنامه", "پروژه", "بلوک", "عنوان", "دسته", "احتمال", "اثر", "امتیاز", "سطح", "مالک", "اقدام پیشگیرانه", "برنامه واکنش", "تاریخ بازبینی", "وضعیت"])
    for r in risks:
        sh.append([r.get("program_title"), r.get("project_title"), r.get("zone_name"), r.get("title"), r.get("category"), r.get("probability"), r.get("impact"), r.get("risk_score"), r.get("risk_level"), r.get("owner"), r.get("mitigation"), r.get("contingency"), r.get("review_date"), r.get("status")])
    style_sheet(sh, [28, 28, 18, 36, 14, 10, 10, 10, 12, 20, 42, 42, 14, 16])

    sh = wb.create_sheet("درخواست تغییر")
    sh.append(["برنامه", "پروژه", "عنوان", "نوع", "فیلد هدف", "درخواست‌کننده", "تاریخ", "اثر زمانی", "اثر هزینه", "مقدار قبلی", "مقدار جدید", "وضعیت", "نظر بررسی"])
    for c in changes:
        sh.append([c.get("program_title"), c.get("project_title"), c.get("title"), c.get("change_type"), c.get("target_field"), c.get("requested_by"), c.get("request_date"), c.get("impact_days"), c.get("impact_cost"), c.get("old_value"), c.get("new_value"), c.get("status"), c.get("review_note")])
    style_sheet(sh, [28, 28, 36, 14, 18, 20, 14, 14, 16, 28, 28, 18, 40])

    sh = wb.create_sheet("هشدارها")
    sh.append(["شدت", "نوع", "عنوان", "بلوک", "سررسید", "پیام"])
    for a in alerts:
        sh.append([a.get("severity"), a.get("type"), a.get("title"), a.get("zone_name"), a.get("due_date"), a.get("message")])
    style_sheet(sh, [12, 20, 36, 20, 14, 60])

    # گانت ماهانه: ۱۲ ستون ماه برای سال مالی منتخب یا سال جاری
    year_text = str(fiscal_year or jalali_year())
    sh = wb.create_sheet("گانت")
    headers = ["کد", "عنوان", "شروع", "پایان", "پیشرفت٪"] + [f"ماه {i}" for i in range(1, 13)]
    sh.append(headers)
    for p in projects:
        row = [p.get("project_code"), p.get("title"), p.get("start_date"), p.get("end_date"), p.get("actual_progress")]
        try:
            start_month = datetime.strptime(str(p.get("start_date"))[:10], "%Y-%m-%d").month
            end_month = datetime.strptime(str(p.get("end_date"))[:10], "%Y-%m-%d").month
        except Exception:
            start_month = end_month = 0
        row.extend(["■" if start_month and start_month <= m <= end_month else "" for m in range(1, 13)])
        sh.append(row)
    style_sheet(sh, [16, 36, 14, 14, 12] + [9] * 12)
    for r in range(2, sh.max_row + 1):
        for c in range(6, 18):
            if sh.cell(r, c).value:
                sh.cell(r, c).fill = PatternFill("solid", fgColor=gold)
                sh.cell(r, c).font = Font(color=gold)
                sh.cell(r, c).alignment = Alignment(horizontal="center")

    wb.save(path)
    return path


def export_project_control_pdf(db, path, fiscal_year=None, zone_id=None):
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from report_generator import _pdf_header_flowables, _pdf_styles
    import report_generator as rg
    from pdf_text_utils import shape_fa

    summary = db.get_project_control_summary(fiscal_year=fiscal_year, zone_id=zone_id)
    programs = db.get_annual_programs(fiscal_year=fiscal_year, zone_id=zone_id)
    projects = db.get_projects(zone_id=zone_id)
    if fiscal_year:
        projects = [p for p in projects if str(p.get("fiscal_year") or "") == str(fiscal_year)]
    risks = db.get_project_risks(zone_id=zone_id, open_only=True)
    changes = db.get_project_change_requests(status="در انتظار بررسی")
    alerts = db.get_project_control_alerts(fiscal_year=fiscal_year, zone_id=zone_id)

    rg._ensure_fonts_registered()
    cell_style = ParagraphStyle("ProjectCellFa", fontName=rg.FONT_NAME, fontSize=7.5, leading=10,
                                alignment=2, wordWrap="RTL")
    head_style = ParagraphStyle("ProjectHeadFa", fontName=rg.FONT_NAME_BOLD, fontSize=8, leading=10,
                                alignment=2, textColor=colors.white, wordWrap="RTL")

    def wrapped_table(headers, rows, widths):
        data = [[Paragraph(shape_fa(str(h)), head_style) for h in headers]]
        for row in rows:
            data.append([Paragraph(shape_fa(str(_safe(v))), cell_style) for v in row])
        table = Table(data, colWidths=widths, repeatRows=1, splitByRow=1)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#13294b")),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#d7dbe3")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f0f2f6")]),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("ALIGN", (0, 0), (-1, -1), "RIGHT"),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("LEFTPADDING", (0, 0), (-1, -1), 3),
            ("RIGHTPADDING", (0, 0), (-1, -1), 3),
        ]))
        return table

    doc = SimpleDocTemplate(path, pagesize=landscape(A4), rightMargin=10*mm, leftMargin=10*mm,
                            topMargin=10*mm, bottomMargin=10*mm)
    story = _pdf_header_flowables("گزارش برنامه عملیاتی و کنترل پروژه")
    _, _, section_style, normal_style = _pdf_styles()
    story.append(Paragraph(shape_fa(f"سال مالی: {fiscal_year or 'همه'}"), normal_style))
    story.append(Spacer(1, 3*mm))
    metrics = [
        ["برنامه", summary["programs_count"], "پروژه", summary["projects_count"], "فعال", summary["active_projects"]],
        ["معوق", summary["overdue_projects"], "میانگین پیشرفت", f"{summary['average_progress']}٪", "تحقق شاخص", f"{summary['indicator_achievement']}٪"],
        ["بودجه", f"{summary['planned_budget']:,.0f}", "هزینه", f"{summary['actual_cost']:,.0f}", "انحراف", f"{summary['cost_variance']:,.0f}"],
        ["ریسک بالا", summary["high_risks"], "تغییر در انتظار", summary["pending_changes"], "هشدار", summary["alerts_count"]],
    ]
    story.append(wrapped_table(["شاخص", "مقدار", "شاخص", "مقدار", "شاخص", "مقدار"], metrics,
                                 [35*mm, 25*mm, 35*mm, 25*mm, 35*mm, 25*mm]))
    story.append(Spacer(1, 5*mm))
    story.append(Paragraph(shape_fa("پروژه‌ها"), section_style))
    project_rows = [[p.get("project_code"), p.get("title"), p.get("program_title"), p.get("zone_name"),
                     p.get("end_date"), f"{float(p.get('actual_progress') or 0):.0f}٪",
                     f"{float(p.get('actual_cost') or 0):,.0f}", p.get("status")]
                    for p in projects]
    if project_rows:
        story.append(wrapped_table(["کد", "عنوان", "برنامه", "بلوک", "پایان", "پیشرفت", "هزینه", "وضعیت"],
                                     project_rows, [25*mm, 52*mm, 45*mm, 32*mm, 25*mm, 22*mm, 30*mm, 28*mm]))
    else:
        story.append(Paragraph(shape_fa("پروژه‌ای ثبت نشده است."), normal_style))
    story.append(PageBreak())
    story.extend(_pdf_header_flowables("ریسک‌ها و تغییرات پروژه"))
    story.append(Paragraph(shape_fa("ریسک‌های باز"), section_style))
    risk_rows = [[r.get("title"), r.get("project_title") or r.get("program_title"), r.get("risk_score"),
                  r.get("risk_level"), r.get("owner"), r.get("review_date"), r.get("status")]
                 for r in risks]
    if risk_rows:
        story.append(wrapped_table(["عنوان", "پروژه/برنامه", "امتیاز", "سطح", "مالک", "بازبینی", "وضعیت"],
                                     risk_rows, [58*mm, 55*mm, 22*mm, 25*mm, 35*mm, 25*mm, 25*mm]))
    story.append(Spacer(1, 5*mm))
    story.append(Paragraph(shape_fa("درخواست‌های تغییر در انتظار"), section_style))
    change_rows = [[c.get("title"), c.get("project_title") or c.get("program_title"), c.get("change_type"),
                    c.get("impact_days"), f"{float(c.get('impact_cost') or 0):,.0f}", c.get("requested_by"), c.get("status")]
                   for c in changes]
    if change_rows:
        story.append(wrapped_table(["عنوان", "پروژه/برنامه", "نوع", "روز", "اثر هزینه", "درخواست‌کننده", "وضعیت"],
                                     change_rows, [55*mm, 55*mm, 25*mm, 18*mm, 30*mm, 35*mm, 28*mm]))
    story.append(Spacer(1, 5*mm))
    story.append(Paragraph(shape_fa("هشدارهای کنترل پروژه"), section_style))
    alert_rows = [[a.get("severity"), a.get("type"), a.get("title"), a.get("zone_name"), a.get("due_date"), a.get("message")]
                  for a in alerts]
    if alert_rows:
        story.append(wrapped_table(["شدت", "نوع", "عنوان", "بلوک", "سررسید", "پیام"], alert_rows,
                                     [22*mm, 35*mm, 50*mm, 30*mm, 25*mm, 85*mm]))
    doc.build(story)
    return path


def export_project_control_powerpoint(db, path, fiscal_year=None, zone_id=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    summary = db.get_project_control_summary(fiscal_year=fiscal_year, zone_id=zone_id)
    programs = db.get_annual_programs(fiscal_year=fiscal_year, zone_id=zone_id)
    projects = db.get_projects(zone_id=zone_id)
    if fiscal_year:
        projects = [p for p in projects if str(p.get("fiscal_year") or "") == str(fiscal_year)]
    risks = db.get_project_risks(zone_id=zone_id, open_only=True)
    alerts = db.get_project_control_alerts(fiscal_year=fiscal_year, zone_id=zone_id)

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    def title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.4))
        tf = box.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(30); p.font.bold = True
        if subtitle:
            p2 = tf.add_paragraph(); p2.text = subtitle; p2.alignment = PP_ALIGN.CENTER; p2.font.size = Pt(18)
        return slide

    def bullet_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7))
        p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(24); p.font.bold = True; p.alignment = PP_ALIGN.RIGHT
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.6), Inches(5.6))
        tf = body.text_frame; tf.word_wrap = True; tf.clear()
        for idx, text in enumerate(bullets):
            para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            para.text = str(text); para.alignment = PP_ALIGN.RIGHT; para.font.size = Pt(18); para.space_after = Pt(8)
        return slide

    title_slide("گزارش برنامه عملیاتی و کنترل پروژه", f"سال مالی: {fiscal_year or 'همه'}")
    bullet_slide("نمای مدیریتی", [
        f"برنامه‌ها: {summary['programs_count']} | پروژه‌ها: {summary['projects_count']}",
        f"پروژه‌های فعال: {summary['active_projects']} | معوق: {summary['overdue_projects']}",
        f"میانگین پیشرفت: {summary['average_progress']}٪ | تحقق شاخص‌ها: {summary['indicator_achievement']}٪",
        f"بودجه: {summary['planned_budget']:,.0f} | هزینه واقعی: {summary['actual_cost']:,.0f}",
        f"ریسک‌های بالا: {summary['high_risks']} | تغییرات در انتظار: {summary['pending_changes']}",
    ])
    bullet_slide("برنامه‌های سالانه", [f"{p.get('fiscal_year')} - {p.get('title')} - پیشرفت {float(p.get('progress_percent') or 0):.0f}٪ - {p.get('status')}" for p in programs[:12]] or ["برنامه‌ای ثبت نشده است."])
    for chunk_start in range(0, min(len(projects), 30), 10):
        chunk = projects[chunk_start:chunk_start + 10]
        bullet_slide(f"پروژه‌ها {chunk_start + 1} تا {chunk_start + len(chunk)}", [
            f"{p.get('project_code')} | {p.get('title')} | پیشرفت {float(p.get('actual_progress') or 0):.0f}٪ | {p.get('status')}" for p in chunk
        ])
    bullet_slide("ریسک‌های مهم", [f"{r.get('risk_level')} ({r.get('risk_score')}) - {r.get('title')} - مالک: {r.get('owner') or 'تعیین نشده'}" for r in risks[:12]] or ["ریسک بازی ثبت نشده است."])
    bullet_slide("هشدارهای کنترل پروژه", [f"{a.get('severity')} - {a.get('title')}: {a.get('message')}" for a in alerts[:12]] or ["هشدار فعالی وجود ندارد."])
    prs.save(path)
    return path


def build_project_control_preview_html(db, fiscal_year=None, zone_id=None):
    """پیش‌نمایش RTL پایدار برای گزارش کنترل پروژه."""
    from html import escape
    summary = db.get_project_control_summary(fiscal_year=fiscal_year, zone_id=zone_id)
    programs = db.get_annual_programs(fiscal_year=fiscal_year, zone_id=zone_id)
    projects = db.get_projects(zone_id=zone_id)
    if fiscal_year:
        projects = [p for p in projects if str(p.get("fiscal_year") or "") == str(fiscal_year)]
    risks = db.get_project_risks(zone_id=zone_id, open_only=True)
    alerts = db.get_project_control_alerts(fiscal_year=fiscal_year, zone_id=zone_id)

    def td(v):
        return f"<td>{escape(convert_dates_in_text(str(_safe(v))))}</td>"

    project_rows = "".join(
        "<tr>" + td(p.get("project_code")) + td(p.get("title")) + td(p.get("program_title")) +
        td(p.get("zone_name")) + td(f"{float(p.get('actual_progress') or 0):.0f}٪") +
        td(p.get("end_date")) + td(p.get("status")) + "</tr>" for p in projects
    ) or "<tr><td colspan='7'>پروژه‌ای ثبت نشده است.</td></tr>"
    risk_rows = "".join(
        "<tr>" + td(r.get("risk_level")) + td(r.get("risk_score")) + td(r.get("title")) +
        td(r.get("project_title") or r.get("program_title")) + td(r.get("owner")) + td(r.get("status")) + "</tr>"
        for r in risks
    ) or "<tr><td colspan='6'>ریسک بازی ثبت نشده است.</td></tr>"
    alert_rows = "".join(
        "<tr>" + td(a.get("severity")) + td(a.get("type")) + td(a.get("title")) +
        td(a.get("zone_name")) + td(a.get("due_date")) + td(a.get("message")) + "</tr>" for a in alerts
    ) or "<tr><td colspan='6'>هشدار فعالی وجود ندارد.</td></tr>"
    metrics = [
        ("برنامه‌ها", summary["programs_count"]), ("پروژه‌ها", summary["projects_count"]),
        ("پروژه فعال", summary["active_projects"]), ("پروژه معوق", summary["overdue_projects"]),
        ("میانگین پیشرفت", f"{summary['average_progress']}٪"),
        ("تحقق شاخص‌ها", f"{summary['indicator_achievement']}٪"),
        ("ریسک بالا", summary["high_risks"]), ("تغییر در انتظار", summary["pending_changes"]),
    ]
    cards = "".join(f"<div class='card'><b>{escape(k)}</b><span>{escape(str(v))}</span></div>" for k, v in metrics)
    return f"""<!doctype html><html lang='fa' dir='rtl'><head><meta charset='utf-8'>
    <style>
    body{{font-family:Tahoma,Arial,sans-serif;background:#f3f5f8;color:#172033;margin:0;padding:24px}}
    .page{{background:white;max-width:1150px;margin:auto;padding:28px;border-radius:12px;box-shadow:0 4px 20px #0001}}
    h1{{color:#13294b;border-bottom:3px solid #c9a227;padding-bottom:12px}} h2{{color:#13294b;margin-top:28px}}
    .cards{{display:grid;grid-template-columns:repeat(4,1fr);gap:10px}} .card{{background:#edf2f8;border:1px solid #d7dfe9;border-radius:9px;padding:14px}}
    .card span{{display:block;font-size:22px;font-weight:bold;color:#13294b;margin-top:6px}}
    table{{width:100%;border-collapse:collapse;margin-top:10px;font-size:13px}} th{{background:#13294b;color:white;padding:9px}} td{{padding:8px;border:1px solid #d8dee7;vertical-align:top}} tr:nth-child(even){{background:#f4f6f9}}
    </style></head><body><div class='page'><h1>گزارش برنامه عملیاتی و کنترل پروژه</h1>
    <p>سال مالی: {escape(str(fiscal_year or 'همه'))}</p><div class='cards'>{cards}</div>
    <h2>پروژه‌ها</h2><table><tr><th>کد</th><th>عنوان</th><th>برنامه</th><th>بلوک</th><th>پیشرفت</th><th>پایان</th><th>وضعیت</th></tr>{project_rows}</table>
    <h2>ریسک‌های باز</h2><table><tr><th>سطح</th><th>امتیاز</th><th>عنوان</th><th>پروژه/برنامه</th><th>مالک</th><th>وضعیت</th></tr>{risk_rows}</table>
    <h2>هشدارها</h2><table><tr><th>شدت</th><th>نوع</th><th>عنوان</th><th>بلوک</th><th>سررسید</th><th>پیام</th></tr>{alert_rows}</table>
    </div></body></html>"""
