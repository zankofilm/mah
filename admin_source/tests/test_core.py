# -*- coding: utf-8 -*-
import hashlib
import json
import os
import sqlite3
import tempfile
import unittest
import zipfile
from unittest.mock import patch

from database import Database
from access_control import has_permission, role_title
from geometry_utils import (clip_polyline_to_polygon, point_in_polygon, polygon_metrics,
                            polygon_self_intersects, polygons_overlap, validate_polygon)
from osm_fetcher import fetch_osm_data
from zone_snapshot_service import export_zone_snapshot_png, refresh_zone_snapshot
from report_generator import (generate_block_full_report_excel, generate_block_full_report_pdf, generate_block_full_report_pptx,
                              generate_correspondence_report_pdf, generate_correspondence_report_excel,
                              generate_correspondence_report_pptx)
from block_report_preview import build_block_full_report_preview_html
from report_preview_html import build_correspondence_report_preview_html
from document_service import generate_document_from_template, render_template_text
from management_calendar_reports import export_management_brief_excel, export_management_brief_pdf
from project_control_reports import export_project_control_excel, export_project_control_pdf, export_project_control_powerpoint
from contracts_satisfaction_reports import export_contract_management_excel, export_contract_management_pdf, export_contract_management_powerpoint


class GeometryTests(unittest.TestCase):
    def test_point_and_clipping(self):
        polygon = [(0, 0), (0, 10), (10, 10), (10, 0)]
        self.assertTrue(point_in_polygon(5, 5, polygon))
        self.assertTrue(point_in_polygon(0, 5, polygon))
        self.assertFalse(point_in_polygon(12, 5, polygon))
        fragments = clip_polyline_to_polygon([(5, -5), (5, 15)], polygon)
        self.assertEqual(len(fragments), 1)
        self.assertAlmostEqual(fragments[0][0][1], 0.0, places=7)
        self.assertAlmostEqual(fragments[0][-1][1], 10.0, places=7)

    def test_polygon_metrics(self):
        area, perimeter = polygon_metrics([(34.8, 46.49), (34.8, 46.50), (34.81, 46.50), (34.81, 46.49)])
        self.assertGreater(area, 900000)
        self.assertGreater(perimeter, 3500)

    def test_polygon_metrics_are_stable_and_orientation_independent(self):
        polygon = [
            (34.8000, 46.4900),
            (34.8000, 46.5000),
            (34.8100, 46.5000),
            (34.8100, 46.4900),
        ]
        area, perimeter = polygon_metrics(polygon)
        reversed_area, reversed_perimeter = polygon_metrics(list(reversed(polygon)))
        closed_area, closed_perimeter = polygon_metrics(polygon + [polygon[0]])
        self.assertAlmostEqual(area, reversed_area, places=6)
        self.assertAlmostEqual(perimeter, reversed_perimeter, places=6)
        self.assertAlmostEqual(area, closed_area, places=6)
        self.assertAlmostEqual(perimeter, closed_perimeter, places=6)
        self.assertAlmostEqual(area, 1_015_120.30, delta=2.0)

    def test_polygon_validation_and_overlap(self):
        valid = [(0, 0), (0, 4), (4, 4), (4, 0)]
        bow_tie = [(0, 0), (4, 4), (0, 4), (4, 0)]
        self.assertFalse(polygon_self_intersects(valid))
        self.assertTrue(polygon_self_intersects(bow_tie))
        ok, message = validate_polygon(valid, minimum_area_m2=0)
        self.assertTrue(ok, message)
        self.assertGreater(polygon_metrics(valid)[0], 0)
        ok, message = validate_polygon(bow_tie, minimum_area_m2=0)
        self.assertFalse(ok)
        self.assertTrue(message)
        self.assertTrue(polygons_overlap(valid, [(3, 3), (3, 6), (6, 6), (6, 3)]))
        self.assertFalse(polygons_overlap(valid, [(5, 5), (5, 6), (6, 6), (6, 5)]))


class DatabaseTests(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.path = os.path.join(self.tmp.name, "app.db")
        self.db = Database(self.path)

    def tearDown(self):
        self.db.close()
        self.tmp.cleanup()

    def test_mosques_and_zone_assignment(self):
        self.assertEqual(len(self.db.get_mosques()), 24)
        mosque = self.db.get_mosques()[0]
        d = 0.0005
        zone_id = self.db.create_zone("مسجد تست", [
            (mosque["lat"] - d, mosque["lon"] - d),
            (mosque["lat"] - d, mosque["lon"] + d),
            (mosque["lat"] + d, mosque["lon"] + d),
            (mosque["lat"] + d, mosque["lon"] - d),
        ])
        assigned = self.db.get_mosques(zone_id=zone_id)
        self.assertIn(mosque["id"], {m["id"] for m in assigned})
        zone = self.db.get_zone(zone_id)
        self.assertGreater(zone["area_m2"], 0)
        self.assertGreater(zone["perimeter_m"], 0)

    def test_city_and_block_area_summary(self):
        city = [
            (34.8000, 46.4900),
            (34.8000, 46.5100),
            (34.8100, 46.5100),
            (34.8100, 46.4900),
        ]
        self.db.save_boundary(city)
        self.db.create_zone("نیمه غربی", [
            (34.8000, 46.4900),
            (34.8000, 46.5000),
            (34.8100, 46.5000),
            (34.8100, 46.4900),
        ])
        self.db.create_zone("نیمه شرقی", [
            (34.8000, 46.5000),
            (34.8000, 46.5100),
            (34.8100, 46.5100),
            (34.8100, 46.5000),
        ])
        summary = self.db.get_area_summary()
        self.assertEqual(summary["zone_count"], 2)
        self.assertGreater(summary["city_area_m2"], 0)
        self.assertAlmostEqual(summary["coverage_percent"], 100.0, places=5)
        self.assertAlmostEqual(summary["difference_m2"], 0.0, places=3)

    def test_atomic_osm_replace_and_segments(self):
        zone_id = self.db.create_zone("ب", [(0, 0), (0, 2), (2, 2), (2, 0)])
        streets = [
            {"osm_id": 10, "segment_index": 0, "name": "الف", "highway_type": "residential", "is_unnamed": 0, "geometry": [(0, 0), (1, 1)]},
            {"osm_id": 10, "segment_index": 1, "name": "الف", "highway_type": "residential", "is_unnamed": 0, "geometry": [(1, 1), (2, 2)]},
        ]
        self.db.replace_osm_data(zone_id, streets=streets, places=[], replace_streets=True, replace_places=True)
        saved = self.db.get_streets(zone_id)
        self.assertEqual(len(saved), 2)
        self.assertEqual({s["segment_index"] for s in saved}, {0, 1})
        # جایگزینی فقط اماکن، خیابان‌ها را حفظ می‌کند
        self.db.replace_osm_data(zone_id, streets=[], places=[], replace_streets=False, replace_places=True)
        self.assertEqual(len(self.db.get_streets(zone_id)), 2)

    def test_places_view_includes_zone_mosques_without_duplication(self):
        mosque = self.db.get_mosques()[0]
        d = 0.0005
        zone_id = self.db.create_zone("نمای اماکن", [
            (mosque["lat"] - d, mosque["lon"] - d),
            (mosque["lat"] - d, mosque["lon"] + d),
            (mosque["lat"] + d, mosque["lon"] + d),
            (mosque["lat"] + d, mosque["lon"] - d),
        ])
        self.db.replace_osm_data(
            zone_id,
            streets=[],
            places=[{
                "osm_id": 100, "name": "مدرسه تست", "category": "آموزشی",
                "subtype": "مدرسه", "lat": mosque["lat"], "lon": mosque["lon"], "address": ""
            }],
            replace_streets=False, replace_places=True,
        )
        rows = self.db.get_places_with_mosques(zone_id)
        mosque_rows = [row for row in rows if row.get("record_type") == "mosque"]
        place_rows = [row for row in rows if row.get("record_type") == "place"]
        self.assertEqual(len(place_rows), 1)
        self.assertTrue(any(row["mosque_id"] == mosque["id"] for row in mosque_rows))
        self.assertEqual(self.db.conn.execute("SELECT COUNT(*) FROM places WHERE zone_id=?", (zone_id,)).fetchone()[0], 1)

    def test_mosque_can_be_meeting_place(self):
        mosque = self.db.get_mosques()[0]
        d = 0.0005
        zone_id = self.db.create_zone("جلسه مسجد", [
            (mosque["lat"] - d, mosque["lon"] - d),
            (mosque["lat"] - d, mosque["lon"] + d),
            (mosque["lat"] + d, mosque["lon"] + d),
            (mosque["lat"] + d, mosque["lon"] - d),
        ])
        self.assertEqual(self.db.get_zone(zone_id)["status"], "ناقص")
        self.db.replace_osm_data(
            zone_id,
            streets=[{"osm_id": 1, "segment_index": 0, "name": "معبر", "highway_type": "residential", "geometry": [(0, 0), (1, 1)]}],
            places=[], replace_streets=True, replace_places=False,
        )
        self.assertEqual(self.db.get_zone(zone_id)["status"], "در حال تکمیل")
        self.db.set_zone_meeting_place(
            zone_id, None, mosque["name"], "آدرس تست", mosque["lat"], mosque["lon"],
            source_type="mosque", source_id=mosque["id"],
        )
        self.assertEqual(self.db.get_zone(zone_id)["status"], "کامل")
        meeting = self.db.get_zone_meeting_place(zone_id)
        self.assertEqual(meeting["source_type"], "mosque")
        self.assertEqual(meeting["source_id"], mosque["id"])
        self.assertIsNone(meeting["place_id"])

    def test_zone_delete_cascades_related_records(self):
        zone_id = self.db.create_zone("حذف آبشاری", [(0, 0), (0, 1), (1, 1), (1, 0)])
        cur = self.db.conn.cursor()
        cur.execute("INSERT INTO council_members (zone_id, first_name, last_name) VALUES (?, 'الف', 'ب')", (zone_id,))
        cur.execute("INSERT INTO priority_requests (zone_id, description) VALUES (?, 'تست')", (zone_id,))
        request_id = cur.lastrowid
        cur.execute("INSERT INTO request_actions (request_id, action_description) VALUES (?, 'پیگیری')", (request_id,))
        self.db.conn.commit()
        self.db.add_field_visit(zone_id, officer_name="کارشناس حذف")
        self.db.add_citizen_request(zone_id, "درخواست حذف")
        self.db.delete_zone(zone_id)
        self.assertEqual(self.db.conn.execute("SELECT COUNT(*) FROM council_members WHERE zone_id=?", (zone_id,)).fetchone()[0], 0)
        self.assertEqual(self.db.conn.execute("SELECT COUNT(*) FROM priority_requests WHERE zone_id=?", (zone_id,)).fetchone()[0], 0)
        self.assertEqual(self.db.conn.execute("SELECT COUNT(*) FROM request_actions WHERE request_id=?", (request_id,)).fetchone()[0], 0)
        self.assertEqual(self.db.conn.execute("SELECT COUNT(*) FROM zone_snapshots WHERE zone_id=?", (zone_id,)).fetchone()[0], 0)
        self.assertEqual(self.db.conn.execute("SELECT COUNT(*) FROM field_visits WHERE zone_id=?", (zone_id,)).fetchone()[0], 0)
        self.assertEqual(self.db.conn.execute("SELECT COUNT(*) FROM citizen_requests WHERE zone_id=?", (zone_id,)).fetchone()[0], 0)
        queued_deletes = self.db.conn.execute(
            "SELECT COUNT(*) FROM offline_sync_queue WHERE operation='delete'"
        ).fetchone()[0]
        self.assertGreaterEqual(queued_deletes, 2)

    def test_password_legacy_migration(self):
        legacy = hashlib.sha256("secret123".encode("utf-8")).hexdigest()
        self.db.conn.execute("UPDATE admin_settings SET username='u', password_hash=? WHERE id=1", (legacy,))
        self.db.conn.commit()
        self.assertTrue(self.db.verify_admin_login("u", "secret123"))
        new_hash = self.db.conn.execute("SELECT password_hash FROM admin_settings WHERE id=1").fetchone()[0]
        self.assertTrue(new_hash.startswith("pbkdf2_sha256$"))

    def test_backup_validation(self):
        backup_path = os.path.join(self.tmp.name, "backup.db")
        self.db.create_backup(backup_path)
        ok, _ = Database.validate_database_file(backup_path)
        self.assertTrue(ok)
        self.assertTrue(os.path.getsize(backup_path) > 0)
        ok, _ = Database.validate_database_file(self.path)
        self.assertTrue(ok)
        invalid = os.path.join(self.tmp.name, "bad.db")
        sqlite3.connect(invalid).close()
        ok, message = Database.validate_database_file(invalid)
        self.assertFalse(ok)
        self.assertTrue(message)

    def test_neighborhood_management_workflow(self):
        zone_id = self.db.create_zone(
            "مدیریت محله تست",
            [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
        )
        profile = self.db.save_zone_profile(
            zone_id, residential_buildings=25, residential_units=120, occupied_units=108,
            vacant_units=12, estimated_households=108, field_households=111,
            approved_households=110, estimated_population=374, average_household_size=3.4,
            elderly_count=30, children_count=75, disabled_count=4, vulnerable_households=12,
            female_headed_households=7, estimation_method="روش ترکیبی", confidence_level="زیاد",
            notes="آزمون پرونده بلوک",
        )
        self.assertEqual(profile["approved_households"], 110)
        self.assertEqual(profile["estimated_population"], 374)

        issue_id = self.db.add_neighborhood_issue(
            zone_id, "روشنایی ناکافی کوچه", category="روشنایی", description="نیازمند نصب چراغ",
            related_office="اداره برق", urgency=5, severity=4, affected_households=80,
            safety_risk=4, status="تأییدشده", due_date="2026-08-01",
        )
        issue = self.db.get_neighborhood_issue(issue_id)
        self.assertIn(issue["priority_level"], ("فوری", "بحرانی"))
        self.assertGreater(issue["priority_score"], 60)

        action_id = self.db.add_neighborhood_action(
            zone_id, "نصب روشنایی جدید", issue_id=issue_id, responsible_office="اداره برق",
            planned_start="2026-07-20", planned_end="2026-08-01", progress_percent=25,
            estimated_cost=500000000, status="در حال اجرا",
        )
        self.assertEqual(self.db.get_neighborhood_issue(issue_id)["status"], "در حال اجرا")
        self.db.update_neighborhood_action(action_id, progress_percent=100, status="تکمیل‌شده")
        self.assertEqual(self.db.get_neighborhood_issue(issue_id)["status"], "مختومه")

        meeting_id = self.db.add_neighborhood_meeting(
            zone_id, "جلسه شورای محله", meeting_date="2026-07-21", start_time="18:00",
            place_name="مسجد محله", agenda="بررسی روشنایی", attendees="اعضای شورا",
            minutes_text="نصب روشنایی تصویب شد",
        )
        resolution_id = self.db.add_neighborhood_resolution(
            meeting_id, zone_id, "پیگیری نصب چراغ", responsible_office="اداره برق",
            due_date="2026-08-01", linked_issue_id=issue_id, linked_action_id=action_id,
        )
        summary = self.db.get_neighborhood_summary(zone_id)
        self.assertEqual(summary["profile"]["approved_households"], 110)
        self.assertEqual(summary["meetings_total"], 1)
        self.assertEqual(summary["resolutions_pending"], 1)
        self.assertEqual(self.db.get_neighborhood_resolution(resolution_id)["linked_action_id"], action_id)

        pdf_path = os.path.join(self.tmp.name, "neighborhood.pdf")
        xlsx_path = os.path.join(self.tmp.name, "neighborhood.xlsx")
        pptx_path = os.path.join(self.tmp.name, "neighborhood.pptx")
        generate_block_full_report_pdf(self.db, zone_id, pdf_path)
        generate_block_full_report_excel(self.db, zone_id, xlsx_path)
        generate_block_full_report_pptx(self.db, zone_id, pptx_path)
        self.assertGreater(os.path.getsize(pdf_path), 1000)
        self.assertGreater(os.path.getsize(xlsx_path), 1000)
        self.assertGreater(os.path.getsize(pptx_path), 1000)

    def test_management_control_budget_alerts_and_performance(self):
        zone_id = self.db.create_zone(
            "کنترل مدیریتی",
            [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
        )
        self.db.save_zone_profile(
            zone_id, residential_buildings=10, residential_units=50, occupied_units=45, vacant_units=5,
            estimated_households=45, field_households=44, approved_households=44, estimated_population=145,
            average_household_size=3.3, elderly_count=8, children_count=30, disabled_count=1,
            vulnerable_households=2, female_headed_households=3, estimation_method="روش ترکیبی",
            confidence_level="متوسط", notes="",
        )
        issue_id = self.db.add_neighborhood_issue(
            zone_id, "آسفالت معبر", related_office="شهرداری", urgency=5, severity=4,
            affected_households=40, safety_risk=4, due_date="2020-01-01",
        )
        action_id = self.db.add_neighborhood_action(
            zone_id, "اجرای آسفالت", issue_id=issue_id, responsible_office="شهرداری",
            planned_start="2020-01-01", planned_end="2020-02-01", progress_percent=30,
            status="در حال اجرا",
        )
        agency_id = self.db.add_management_agency(
            name="شهرداری", category="خدمات شهری", contact_person="مسئول خدمات", phone="123",
        )
        budget_id = self.db.add_neighborhood_budget(
            zone_id, "اعتبار آسفالت", action_id=action_id, fiscal_year="۱۴۰۵",
            funding_source="اعتبارات شهری", approved_amount=1000, allocated_amount=800,
            spent_amount=900, status="در حال هزینه",
        )
        self.assertEqual(self.db.get_neighborhood_budget(budget_id)["action_id"], action_id)
        self.assertEqual(self.db.get_budget_summary(zone_id)["overrun_count"], 1)
        alerts = self.db.get_management_alerts(zone_id)
        self.assertTrue(any(a["category"] == "بودجه" for a in alerts))
        self.assertTrue(any(a["category"] in ("پروژه", "سررسید") for a in alerts))
        self.db.acknowledge_management_alert(alerts[0]["key"])
        self.assertFalse(any(a["key"] == alerts[0]["key"] for a in self.db.get_management_alerts(zone_id)))
        performance = self.db.get_zone_performance(zone_id)
        self.assertGreaterEqual(performance["total_score"], 0)
        self.assertLessEqual(performance["total_score"], 100)
        agency = next(x for x in self.db.get_agency_performance() if x["id"] == agency_id)
        self.assertGreaterEqual(agency["assigned"], 2)
        self.assertGreaterEqual(agency["overdue"], 1)
        stats = self.db.get_system_stats()
        self.assertEqual(stats["budget_overruns_count"], 1)
        self.assertEqual(stats["agencies_count"], 1)

    def test_field_operations_citizen_requests_sync_and_analysis(self):
        zone_id = self.db.create_zone(
            "بلوک عملیات میدانی",
            [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
        )
        visit_id = self.db.add_field_visit(
            zone_id, visit_date="2026-07-20", start_time="09:30", officer_name="کارشناس تست",
            visit_type="شمارش خانوار", location_text="خیابان نمونه", lat=34.805, lon=46.495,
            buildings_count=12, households_count=48, observation="بازدید انجام شد",
            followup_required=True, status="نیازمند پیگیری",
        )
        self.assertEqual(self.db.get_field_visit(visit_id)["households_count"], 48)

        request_id = self.db.add_citizen_request(
            zone_id, "روشنایی معبر", category="روشنایی", description="چراغ معبر خاموش است",
            citizen_name="شهروند تست", mobile="09120000000", location_text="کوچه نمونه",
            lat=34.806, lon=46.496, urgency=5, assigned_office="اداره برق", source="فرم میدانی",
        )
        request = self.db.get_citizen_request(request_id)
        self.assertTrue(request["tracking_code"].startswith("JR-"))
        issue_id = self.db.convert_citizen_request_to_issue(request_id)
        self.assertEqual(self.db.get_citizen_request(request_id)["linked_issue_id"], issue_id)

        alerts = self.db.get_management_alerts(zone_id)
        self.assertTrue(any(a["entity_type"] == "citizen_request" for a in alerts))
        self.assertTrue(any(a["entity_type"] == "field_visit" for a in alerts))
        analysis = self.db.get_zone_operational_analysis(zone_id)
        self.assertEqual(analysis["field_visits"], 1)
        self.assertEqual(analysis["citizen_requests"], 1)
        self.assertGreater(analysis["risk_score"], 0)
        geojson = self.db.build_operational_geojson()
        self.assertEqual(geojson["type"], "FeatureCollection")
        self.assertTrue(any(f["geometry"]["type"] == "Point" for f in geojson["features"]))

        pdf_path = os.path.join(self.tmp.name, "operations.pdf")
        xlsx_path = os.path.join(self.tmp.name, "operations.xlsx")
        pptx_path = os.path.join(self.tmp.name, "operations.pptx")
        generate_block_full_report_pdf(self.db, zone_id, pdf_path)
        generate_block_full_report_excel(self.db, zone_id, xlsx_path)
        generate_block_full_report_pptx(self.db, zone_id, pptx_path)
        self.assertGreater(os.path.getsize(pdf_path), 1000)
        self.assertGreater(os.path.getsize(xlsx_path), 1000)
        self.assertGreater(os.path.getsize(pptx_path), 1000)
        from openpyxl import load_workbook
        from pptx import Presentation
        workbook = load_workbook(xlsx_path, read_only=True)
        self.assertIn("بازدیدهای میدانی", workbook.sheetnames)
        self.assertIn("درخواست‌های مردمی", workbook.sheetnames)
        self.assertIn("تحلیل عملیاتی", workbook.sheetnames)
        workbook.close()
        self.assertGreaterEqual(len(Presentation(pptx_path).slides), 13)
        preview = build_block_full_report_preview_html(self.db, zone_id)
        self.assertIn("بازدیدها و برداشت‌های میدانی", preview)
        self.assertIn("درخواست‌ها و گزارش‌های مردمی", preview)
        self.assertIn("تحلیل عملیاتی بلوک", preview)

        package_path = os.path.join(self.tmp.name, "sync.json")
        result = self.db.export_sync_package(package_path, zone_id=zone_id)
        self.assertGreaterEqual(result["count"], 2)
        self.assertTrue(os.path.exists(package_path))

        other_path = os.path.join(self.tmp.name, "other.db")
        other = Database(other_path)
        try:
            other.create_zone("بلوک غیرمرتبط", [(0, 0), (0, 1), (1, 1), (1, 0)])
            target_zone_id = other.create_zone(
                "بلوک عملیات میدانی",
                [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
            )
            counts = other.import_sync_package(package_path)
            self.assertGreaterEqual(counts["inserted"], 2)
            self.assertEqual(len(other.get_field_visits(target_zone_id)), 1)
            self.assertEqual(len(other.get_citizen_requests(target_zone_id)), 1)
            # ورود مجدد بسته نباید رکورد تکراری بسازد.
            other.import_sync_package(package_path)
            self.assertEqual(len(other.get_field_visits(target_zone_id)), 1)
            self.assertEqual(len(other.get_citizen_requests(target_zone_id)), 1)
        finally:
            other.close()

    def test_v63_users_permissions_audit_search_and_daily_backup(self):
        admin = self.db.authenticate_user("admin", "admin123")
        self.assertIsNotNone(admin)
        self.assertEqual(admin["role"], "admin")
        self.assertTrue(self.db.current_user_can("system_settings"))

        manager_id = self.db.create_user(
            "manager1", "مدیر محله آزمایشی", "StrongPass123", role="manager",
            mobile="09120000000", must_change_password=False,
        )
        self.assertTrue(manager_id)
        self.assertTrue(has_permission("manager", "neighborhood"))
        self.assertFalse(has_permission("manager", "system_settings"))
        self.assertEqual(role_title("gis"), "کارشناس GIS")
        with self.assertRaises(ValueError):
            self.db.update_user(admin["id"], role="viewer")

        self.db.set_current_user(None)
        manager = self.db.authenticate_user("manager1", "StrongPass123")
        self.assertIsNotNone(manager)
        self.assertTrue(self.db.current_user_can("reports"))
        self.assertFalse(self.db.current_user_can("system_settings"))

        zone_id = self.db.create_zone(
            "بلوک جستجوی امنیتی",
            [(34.8, 46.49), (34.8, 46.50), (34.81, 46.50), (34.81, 46.49)],
        )
        self.db.save_streets_bulk([{
            "osm_id": 9898, "segment_index": 0, "name": "خیابان امنیت",
            "highway_type": "residential", "geometry": [(34.801, 46.491), (34.809, 46.499)],
        }], zone_id=zone_id)
        issue_id = self.db.add_neighborhood_issue(
            zone_id, "روشنایی میدان امنیت", description="نیازمند پیگیری روشنایی",
            related_office="شهرداری", urgency=4, severity=3,
        )
        results = self.db.global_search("امنیت")
        result_types = {item["entity_type"] for item in results}
        self.assertIn("zone", result_types)
        self.assertIn("street", result_types)
        self.assertIn("issue", result_types)
        self.assertTrue(any(item.get("zone_id") == zone_id for item in results))

        logs = self.db.get_audit_logs(username="manager1", zone_id=zone_id)
        self.assertTrue(any(item["action"] == "create" for item in logs))
        self.assertTrue(any(item.get("actor_username") == "manager1" for item in logs))

        backup_path = self.db.ensure_daily_backup(keep=3)
        self.assertTrue(backup_path and os.path.exists(backup_path))
        registry_before = self.db.list_registered_backups()
        self.assertGreaterEqual(len(registry_before), 1)
        same_path = self.db.ensure_daily_backup(keep=3)
        self.assertEqual(same_path, backup_path)
        self.assertEqual(len(self.db.list_registered_backups()), len(registry_before))
        stats = self.db.get_system_stats()
        self.assertGreaterEqual(stats["active_users_count"], 2)
        self.assertGreaterEqual(stats["healthy_backups_count"], 1)
        self.assertEqual(self.db.get_neighborhood_issue(issue_id)["zone_id"], zone_id)

        locked_id = self.db.create_user(
            "locked1", "کاربر قفل آزمایشی", "AnotherPass123", role="viewer",
            must_change_password=False,
        )
        self.db.set_current_user(None)
        for _ in range(5):
            self.assertIsNone(self.db.authenticate_user("locked1", "wrong-password"))
        self.assertIsNone(self.db.authenticate_user("locked1", "AnotherPass123"))
        locked = self.db.get_user(locked_id)
        self.assertTrue(locked.get("locked_until"))

    def test_v64_correspondence_attachments_workflow_and_reports(self):
        admin = self.db.authenticate_user("admin", "admin123")
        self.assertIsNotNone(admin)
        zone_id = self.db.create_zone(
            "بلوک مکاتبات",
            [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
        )
        staff_id = self.db.create_user(
            "staff64", "کارشناس مکاتبات", "StrongPass456", role="field",
            must_change_password=False,
        )
        letter_id = self.db.add_correspondence_letter(
            "1405/100", "وارده", "درخواست اصلاح روشنایی محله", zone_id=zone_id,
            sender="فرمانداری", recipient="شهرداری", letter_date="2026-07-01",
            received_date="2026-07-02", due_date="2026-07-10", priority="فوری",
            description="نامه آزمایشی برای کنترل کارتابل",
        )
        letter = self.db.get_correspondence_letter(letter_id)
        self.assertEqual(letter["zone_id"], zone_id)
        self.assertEqual(letter["direction"], "وارده")
        with self.assertRaises(ValueError):
            self.db.add_correspondence_letter("1405/100", "وارده", "نامه تکراری")

        source_path = os.path.join(self.tmp.name, "letter.txt")
        with open(source_path, "w", encoding="utf-8") as handle:
            handle.write("attachment test")
        attachment_id = self.db.archive_document_attachment("letter", letter_id, source_path)
        duplicate_id = self.db.archive_document_attachment("letter", letter_id, source_path)
        self.assertEqual(attachment_id, duplicate_id)
        attachments = self.db.get_document_attachments("letter", letter_id)
        self.assertEqual(len(attachments), 1)
        archived_path = attachments[0]["stored_path"]
        self.assertTrue(os.path.exists(archived_path))

        assignment_id = self.db.add_workflow_assignment(
            letter_id, assigned_to_user_id=staff_id, instruction="بررسی و پاسخ",
            due_date="2026-07-11", priority="فوری",
        )
        assignment = self.db.get_workflow_assignment(assignment_id)
        self.assertEqual(assignment["assigned_to_user_id"], staff_id)
        self.assertEqual(self.db.get_correspondence_letter(letter_id)["status"], "ارجاع‌شده")
        notifications = self.db.get_administrative_notifications()
        self.assertTrue(any(item["type"] == "نامه" for item in notifications))
        self.assertTrue(any(item["type"] == "ارجاع" for item in notifications))
        self.db.acknowledge_administrative_notification(notifications[0]["key"])
        self.assertFalse(any(item["key"] == notifications[0]["key"] for item in self.db.get_administrative_notifications()))

        self.db.update_workflow_assignment(assignment_id, status="پاسخ‌داده‌شده", response_text="پاسخ ثبت شد")
        self.assertEqual(self.db.get_workflow_assignment(assignment_id)["status"], "پاسخ‌داده‌شده")
        self.assertEqual(self.db.get_correspondence_letter(letter_id)["status"], "پاسخ‌داده‌شده")
        results = self.db.global_search("روشنایی")
        self.assertTrue(any(item["entity_type"] == "letter" for item in results))
        stats = self.db.get_system_stats()
        self.assertEqual(stats["letters_total"], 1)
        self.assertEqual(stats["attachments_count"], 1)

        pdf_path = os.path.join(self.tmp.name, "correspondence.pdf")
        xlsx_path = os.path.join(self.tmp.name, "correspondence.xlsx")
        pptx_path = os.path.join(self.tmp.name, "correspondence.pptx")
        generate_block_full_report_pdf(self.db, zone_id, pdf_path)
        generate_block_full_report_excel(self.db, zone_id, xlsx_path)
        generate_block_full_report_pptx(self.db, zone_id, pptx_path)
        self.assertGreater(os.path.getsize(pdf_path), 1000)
        from openpyxl import load_workbook
        from pptx import Presentation
        workbook = load_workbook(xlsx_path, read_only=True)
        self.assertIn("مکاتبات اداری", workbook.sheetnames)
        workbook.close()
        self.assertGreaterEqual(len(Presentation(pptx_path).slides), 14)
        preview = build_block_full_report_preview_html(self.db, zone_id)
        self.assertIn("مکاتبات اداری مرتبط", preview)

        letters_pdf = os.path.join(self.tmp.name, "letters.pdf")
        letters_xlsx = os.path.join(self.tmp.name, "letters.xlsx")
        letters_pptx = os.path.join(self.tmp.name, "letters.pptx")
        generate_correspondence_report_pdf(self.db, letters_pdf, zone_id=zone_id)
        generate_correspondence_report_excel(self.db, letters_xlsx, zone_id=zone_id)
        generate_correspondence_report_pptx(self.db, letters_pptx, zone_id=zone_id)
        self.assertGreater(os.path.getsize(letters_pdf), 1000)
        letters_workbook = load_workbook(letters_xlsx, read_only=True)
        self.assertIn("ارجاعات", letters_workbook.sheetnames)
        letters_workbook.close()
        self.assertGreaterEqual(len(Presentation(letters_pptx).slides), 3)
        self.assertIn("دفتر مکاتبات", build_correspondence_report_preview_html(self.db, zone_id))

        archive_path = os.path.join(self.tmp.name, "correspondence_archive.zip")
        archive_result = self.db.export_correspondence_archive(archive_path, zone_id=zone_id)
        self.assertEqual(archive_result["letters"], 1)
        self.assertEqual(archive_result["files"], 1)
        with zipfile.ZipFile(archive_path, "r") as archive:
            names = archive.namelist()
            self.assertIn("correspondence.json", names)
            self.assertTrue(any(name.startswith("files/") for name in names))

        self.assertTrue(self.db.delete_correspondence_letter(letter_id))
        self.assertEqual(len(self.db.get_workflow_assignments(letter_id=letter_id)), 0)
        self.assertFalse(os.path.exists(archived_path))

    def test_v65_approval_templates_and_decision_dashboard(self):
        admin = self.db.authenticate_user("admin", "admin123")
        self.assertIsNotNone(admin)
        manager_id = self.db.create_user(
            "manager65", "مدیر محله آزمایشی", "StrongPass789", role="manager",
            must_change_password=False,
        )
        manager = self.db.get_user(manager_id)
        zone_id = self.db.create_zone(
            "بلوک گردش تأیید",
            [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
        )
        action_id = self.db.add_neighborhood_action(
            zone_id, "بهسازی روشنایی", responsible_office="شهرداری",
            responsible_person="کارشناس پروژه", planned_end="2026-08-15",
            progress_percent=20, status="در حال اجرا",
        )
        self.db.set_current_user(manager)
        approval_id = self.db.create_approval_request(
            "action", action_id, due_date="2026-08-01",
            steps=[
                {"approver_role": "manager", "approver_name": "مدیر محله"},
                {"approver_role": "admin", "approver_name": "مدیر سامانه"},
            ],
        )
        self.assertEqual(self.db.get_neighborhood_action(action_id)["approval_status"], "در انتظار تأیید")
        with self.assertRaises(ValueError):
            self.db.create_approval_request("action", action_id)
        self.assertTrue(self.db.current_user_can_decide_approval(approval_id))
        first = self.db.decide_approval(approval_id, approved=True, comment="تأیید مرحله اول")
        self.assertEqual(first["status"], "در انتظار تأیید")
        self.assertEqual(first["current_step"], 2)
        self.db.set_current_user(admin)
        final = self.db.decide_approval(approval_id, approved=True, comment="تأیید نهایی")
        self.assertEqual(final["status"], "تأییدشده")
        self.assertEqual(self.db.get_neighborhood_action(action_id)["approval_status"], "تأییدشده")
        self.assertEqual(len(final["steps"]), 2)

        templates = self.db.get_document_templates(active_only=True)
        self.assertGreaterEqual(len(templates), 4)
        letter_template = next(item for item in templates if item["template_type"] == "نامه اداری")
        context = self.db.build_document_context(
            zone_id=zone_id, action_id=action_id,
            extra={"subject": "بهسازی روشنایی", "due_date": "2026-08-15"},
        )
        output_path = os.path.join(self.tmp.name, "generated_letter.docx")
        result = generate_document_from_template(
            self.db, letter_template["id"], output_path, context=context, zone_id=zone_id,
            related_entity_type="action", related_entity_id=action_id,
            metadata={"number": "65/1", "creator": "مدیر سامانه"},
        )
        self.assertTrue(os.path.exists(result["path"]))
        self.assertGreater(os.path.getsize(result["path"]), 10000)
        self.assertIn("بهسازی روشنایی", result["subject"])
        self.assertEqual(len(self.db.get_generated_documents(zone_id)), 1)
        self.assertEqual(render_template_text("{zone_name}-{missing}", {"zone_name": "الف"}), "الف-—")

        decision_rows = self.db.get_zone_decision_rows()
        self.assertEqual(len(decision_rows), 1)
        self.assertEqual(decision_rows[0]["zone_id"], zone_id)
        self.assertIn("score", decision_rows[0])
        stats = self.db.get_system_stats()
        self.assertEqual(stats["approvals_total"], 1)
        self.assertGreaterEqual(stats["document_templates_count"], 4)
        self.assertEqual(stats["generated_documents_count"], 1)
        search = self.db.global_search("بهسازی روشنایی")
        self.assertTrue(any(item["entity_type"] in {"action", "generated_document"} for item in search))

    def test_management_calendar_notifications_and_period_brief(self):
        zone_id = self.db.create_zone(
            "تقویم مدیریتی",
            [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
        )
        action_id = self.db.add_neighborhood_action(
            zone_id, "اجرای روشنایی", responsible_person="مسئول پروژه",
            planned_start="2026-07-18", planned_end="2026-07-20",
            progress_percent=35, status="در حال اجرا",
        )
        event_id = self.db.add_management_calendar_event(
            "جلسه بررسی پیشرفت", "2026-07-21", zone_id=zone_id, category="جلسه",
            start_time="18:00", all_day=False, responsible_person="دبیر جلسه",
            location="فرمانداری", priority="مهم", reminder_days=3,
        )
        event = self.db.get_management_calendar_event(event_id)
        self.assertEqual(event["title"], "جلسه بررسی پیشرفت")
        self.assertFalse(event["all_day"])

        items = self.db.get_deadline_calendar_items("2026-07-19", "2026-07-22", zone_id=zone_id)
        self.assertTrue(any(item["source_type"] == "action" and item["source_id"] == action_id for item in items))
        self.assertTrue(any(item["source_type"] == "calendar_event" and item["source_id"] == event_id for item in items))

        generated = self.db.refresh_in_app_notifications(reference_date="2026-07-19", days_ahead=7)
        self.assertGreaterEqual(generated, 2)
        notifications = self.db.get_in_app_notifications()
        self.assertTrue(any(item["source_type"] == "calendar_event" for item in notifications))
        first_id = notifications[0]["id"]
        self.db.mark_notification_read(first_id)
        self.assertTrue(next(x for x in self.db.get_in_app_notifications() if x["id"] == first_id)["is_read"])

        # توجه: actions_created در get_management_period_brief بر اساس created_at
        # واقعی رکورد فیلتر می‌شود (نه planned_start که در بالا تاریخ ثابت گذشته
        # دارد)، پس این بخش باید بازه‌ای شامل «امروز واقعی» بسنجد.
        from datetime import date, timedelta
        today = date.today()
        brief_today = self.db.get_management_period_brief(
            (today - timedelta(days=1)).isoformat(), (today + timedelta(days=1)).isoformat(), zone_id
        )
        self.assertEqual(brief_today["actions_created"], 1)

        # calendar_events بر اساس start_date ثابتی که بالاتر ثبت شد فیلتر می‌شود
        brief_fixed = self.db.get_management_period_brief("2026-07-18", "2026-07-22", zone_id)
        self.assertEqual(brief_fixed["calendar_events"], 1)
        self.assertGreaterEqual(brief_fixed["deadlines_total"], 2)

        excel_path = os.path.join(self.tmp.name, "management_brief.xlsx")
        pdf_path = os.path.join(self.tmp.name, "management_brief.pdf")
        export_management_brief_excel(self.db, "2026-07-18", "2026-07-22", excel_path, zone_id)
        export_management_brief_pdf(self.db, "2026-07-18", "2026-07-22", pdf_path, zone_id)
        self.assertGreater(os.path.getsize(excel_path), 1000)
        self.assertGreater(os.path.getsize(pdf_path), 1000)
        self.assertTrue(any(x["entity_type"] == "calendar_event" for x in self.db.global_search("بررسی پیشرفت")))

        self.db.delete_management_calendar_event(event_id)
        self.assertIsNone(self.db.get_management_calendar_event(event_id))

    def test_v67_operational_program_project_risk_change_and_reports(self):
        zone_id = self.db.create_zone(
            "بلوک کنترل پروژه",
            [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
        )
        program_id = self.db.add_annual_program(
            "1405", "برنامه بهبود روشنایی", strategic_goal="ارتقای ایمنی شبانه",
            zone_id=zone_id, responsible_agency="شهرداری", program_manager="مدیر پروژه",
            start_date="2026-01-01", end_date="2026-12-31", approved_budget=1_000_000_000,
            status="مصوب",
        )
        project_id = self.db.add_project(
            "نصب چراغ معابر", program_id=program_id, zone_id=zone_id,
            responsible_agency="شهرداری", project_manager="کارشناس فنی",
            start_date="2026-01-01", end_date="2026-08-01", planned_budget=800_000_000,
            planned_progress=70, actual_progress=45, status="در حال اجرا", priority="مهم",
        )
        project = self.db.get_project(project_id)
        self.assertTrue(project["project_code"].startswith("PRJ-"))
        milestone_id = self.db.add_project_milestone(
            project_id, "خرید تجهیزات", due_date="2026-06-01", weight=2,
        )
        self.db.add_project_progress_update(
            project_id, "2026-07-01", planned_progress=70, actual_progress=50,
            actual_cost=500_000_000, summary="نیمی از کار انجام شد",
            obstacles="تأخیر تأمین", next_steps="تحویل تجهیزات",
        )
        indicator_id = self.db.add_project_indicator(
            "تعداد چراغ نصب‌شده", program_id=program_id, project_id=project_id,
            unit="عدد", baseline_value=0, target_value=40, actual_value=20,
        )
        indicator = self.db.get_project_indicator(indicator_id)
        self.assertEqual(indicator["achievement_percent"], 50.0)
        risk_id = self.db.add_project_risk(
            "تأخیر تأمین تجهیزات", program_id=program_id, project_id=project_id,
            category="تأمین", probability=5, impact=4, owner="کارپرداز",
            mitigation="خرید زودهنگام", contingency="تأمین‌کننده جایگزین",
            review_date="2026-07-20",
        )
        self.assertEqual(self.db.get_project_risk(risk_id)["risk_level"], "بحرانی")
        change_id = self.db.add_project_change_request(
            "تمدید پایان پروژه", program_id=program_id, project_id=project_id,
            change_type="زمان", target_field="end_date", reason="تأخیر خرید",
            requested_by="مدیر پروژه", impact_days=31,
            old_value="2026-08-01", new_value="2026-09-01",
        )
        reviewed = self.db.review_project_change_request(
            change_id, "تأییدشده", "با تمدید موافقت شد", apply_change=True,
        )
        self.assertEqual(reviewed["status"], "اعمال‌شده")
        self.assertEqual(self.db.get_project(project_id)["end_date"], "2026-09-01")

        summary = self.db.get_project_control_summary(fiscal_year="1405", zone_id=zone_id)
        self.assertEqual(summary["programs_count"], 1)
        self.assertEqual(summary["projects_count"], 1)
        self.assertEqual(summary["high_risks"], 1)
        self.assertGreater(summary["alerts_count"], 0)
        calendar = self.db.get_deadline_calendar_items("2026-01-01", "2026-12-31", zone_id=zone_id)
        self.assertTrue(any(item["source_type"] == "project" for item in calendar))
        self.assertTrue(any(item["source_type"] == "project_milestone" for item in calendar))
        self.assertTrue(any(item["source_type"] == "project_risk" for item in calendar))
        search = self.db.global_search("چراغ معابر")
        self.assertTrue(any(item["entity_type"] == "project" for item in search))
        stats = self.db.get_system_stats()
        self.assertEqual(stats["annual_programs_count"], 1)
        self.assertEqual(stats["portfolio_projects_count"], 1)
        self.assertEqual(stats["high_project_risks"], 1)
        self.assertTrue(has_permission("manager", "project_control"))
        self.assertTrue(has_permission("viewer", "project_control"))

        pdf_path = os.path.join(self.tmp.name, "project_control.pdf")
        xlsx_path = os.path.join(self.tmp.name, "project_control.xlsx")
        pptx_path = os.path.join(self.tmp.name, "project_control.pptx")
        export_project_control_pdf(self.db, pdf_path, fiscal_year="1405", zone_id=zone_id)
        export_project_control_excel(self.db, xlsx_path, fiscal_year="1405", zone_id=zone_id)
        export_project_control_powerpoint(self.db, pptx_path, fiscal_year="1405", zone_id=zone_id)
        self.assertGreater(os.path.getsize(pdf_path), 1000)
        self.assertGreater(os.path.getsize(xlsx_path), 1000)
        self.assertGreater(os.path.getsize(pptx_path), 1000)
        from openpyxl import load_workbook
        from pptx import Presentation
        workbook = load_workbook(xlsx_path, read_only=True)
        self.assertIn("گانت", workbook.sheetnames)
        self.assertIn("ریسک‌ها", workbook.sheetnames)
        workbook.close()
        self.assertGreaterEqual(len(Presentation(pptx_path).slides), 5)

        self.assertTrue(self.db.delete_project(project_id))
        self.assertIsNone(self.db.get_project_milestone(milestone_id))
        self.assertIsNone(self.db.get_project_indicator(indicator_id))
        self.assertIsNone(self.db.get_project_risk(risk_id))

    def test_zone_snapshot_is_stored_and_regenerated(self):
        zone_id = self.db.create_zone(
            "تصویر تست",
            [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
        )
        first = self.db.get_zone_snapshot(zone_id)
        self.assertIsNotNone(first)
        self.assertEqual(first["render_status"], "ready")
        self.assertTrue(first["png_data"].startswith(b"\x89PNG"))
        self.assertIn("<svg", first["svg_text"])
        first_version = first["version"]

        self.db.replace_osm_data(
            zone_id,
            streets=[{
                "osm_id": 777, "segment_index": 0, "name": "خیابان آزمایشی",
                "highway_type": "residential", "geometry": [(34.801, 46.491), (34.809, 46.499)],
            }],
            places=[], replace_streets=True, replace_places=True,
        )
        second = self.db.get_zone_snapshot(zone_id)
        self.assertGreater(second["version"], first_version)
        self.assertNotEqual(second["content_hash"], first["content_hash"])

        out_path = os.path.join(self.tmp.name, "zone.png")
        self.assertTrue(export_zone_snapshot_png(self.db, zone_id, out_path))
        self.assertGreater(os.path.getsize(out_path), 1000)

        pdf_path = os.path.join(self.tmp.name, "zone.pdf")
        xlsx_path = os.path.join(self.tmp.name, "zone.xlsx")
        pptx_path = os.path.join(self.tmp.name, "zone.pptx")
        generate_block_full_report_pdf(self.db, zone_id, pdf_path)
        generate_block_full_report_excel(self.db, zone_id, xlsx_path)
        generate_block_full_report_pptx(self.db, zone_id, pptx_path)
        self.assertGreater(os.path.getsize(pdf_path), 1000)
        self.assertGreater(os.path.getsize(xlsx_path), 1000)
        self.assertGreater(os.path.getsize(pptx_path), 1000)


    def test_v68_contracts_payments_evaluation_satisfaction_and_reports(self):
        zone_id = self.db.create_zone(
            "بلوک قرارداد",
            [(34.800, 46.490), (34.800, 46.500), (34.810, 46.500), (34.810, 46.490)],
        )
        program_id = self.db.add_annual_program("1405", "برنامه عمرانی", zone_id=zone_id, status="مصوب")
        project_id = self.db.add_project(
            "بهسازی معابر", program_id=program_id, zone_id=zone_id,
            start_date="2026-01-01", end_date="2026-08-31",
            planned_budget=2_000_000_000, status="در حال اجرا",
        )
        contractor_id = self.db.add_contractor(
            "شرکت عمران نمونه", national_id="14001234567", manager_name="مدیر نمونه",
            phone="08300000000", specialty="راه و ابنیه",
        )
        contract_id = self.db.add_contract(
            "C-1405-01", "اجرای بهسازی معابر", contractor_id, project_id=project_id,
            contract_date="2026-01-10", start_date="2026-01-15", end_date="2026-07-01",
            amount=1_500_000_000, guarantee_amount=150_000_000, status="فعال",
        )
        payment_id = self.db.add_contract_payment(
            contract_id, statement_no="1", gross_amount=500_000_000,
            deductions=50_000_000, approved_amount=450_000_000,
            paid_amount=300_000_000, approval_date="2026-06-20",
            status="پرداخت جزئی",
        )
        payment = self.db.get_contract_payment(payment_id)
        self.assertEqual(payment["net_amount"], 450_000_000)
        evaluation_id = self.db.add_contractor_evaluation(
            contract_id, "2026-06-30", quality_score=4.5, schedule_score=4,
            safety_score=5, cooperation_score=4, documentation_score=3.5,
            evaluator="ناظر پروژه",
        )
        self.assertGreater(self.db.get_contractor(contractor_id)["average_score"], 80)
        survey_id = self.db.add_satisfaction_survey(
            "2026-07-15", zone_id=zone_id, project_id=project_id, respondents=25,
            problem_resolved_percent=40, quality_score=2, speed_score=2,
            communication_score=2.5, overall_score=2, reopen_recommended=True,
            recorded_by="کارشناس اجتماعی", comments="نیازمند اصلاح کیفیت اجرا",
        )
        participation_id = self.db.add_community_participation(
            zone_id, "همکاری اهالی در پاکسازی", participation_type="داوطلبانه",
            organization_name="شورای محله", volunteers_count=30,
            cash_value=50_000_000, noncash_value=20_000_000,
            start_date="2026-07-01", status="تکمیل‌شده",
        )
        summary = self.db.get_contract_management_summary(zone_id)
        self.assertEqual(summary["contracts_count"], 1)
        self.assertEqual(summary["paid_amount"], 300_000_000)
        self.assertEqual(summary["volunteers_count"], 30)
        self.assertLess(summary["average_satisfaction"], 50)
        alerts = self.db.get_contract_management_alerts(zone_id)
        self.assertTrue(any(x["type"] == "قرارداد معوق" for x in alerts))
        self.assertTrue(any(x["type"] == "مطالبه پرداخت‌نشده" for x in alerts))
        self.assertTrue(any(x["type"] == "رضایت پایین" for x in alerts))
        calendar = self.db.get_deadline_calendar_items("2026-01-01", "2026-12-31", zone_id=zone_id)
        self.assertTrue(any(x["source_type"] == "contract" for x in calendar))
        self.assertTrue(any(x["entity_type"] == "contract" for x in self.db.global_search("C-1405-01")))
        self.assertTrue(any(x["entity_type"] == "contractor" for x in self.db.global_search("عمران نمونه")))
        stats = self.db.get_system_stats()
        self.assertEqual(stats["contracts_count"], 1)
        self.assertTrue(has_permission("manager", "contracts"))
        self.assertTrue(has_permission("viewer", "contracts"))

        pdf_path = os.path.join(self.tmp.name, "contracts.pdf")
        xlsx_path = os.path.join(self.tmp.name, "contracts.xlsx")
        pptx_path = os.path.join(self.tmp.name, "contracts.pptx")
        export_contract_management_pdf(self.db, pdf_path, zone_id)
        export_contract_management_excel(self.db, xlsx_path, zone_id)
        export_contract_management_powerpoint(self.db, pptx_path, zone_id)
        self.assertGreater(os.path.getsize(pdf_path), 1000)
        self.assertGreater(os.path.getsize(xlsx_path), 1000)
        self.assertGreater(os.path.getsize(pptx_path), 1000)

        self.db.delete_contract(contract_id)
        self.assertEqual(self.db.get_contract_payments(contract_id), [])
        self.assertEqual(self.db.get_contractor_evaluations(contract_id=contract_id), [])
        self.db.delete_contractor(contractor_id)
        self.assertIsNone(self.db.get_contractor(contractor_id))
        self.assertTrue(self.db.delete_satisfaction_survey(survey_id))
        self.assertTrue(self.db.delete_community_participation(participation_id))


class OSMProcessingTests(unittest.TestCase):
    def test_clipping_and_independent_sections(self):
        polygon = [(0, 0), (0, 10), (10, 10), (10, 0)]
        streets_payload = {
            "elements": [{
                "type": "way", "id": 99,
                "tags": {"highway": "residential"},
                "geometry": [{"lat": 5, "lon": -2}, {"lat": 5, "lon": 12}],
            }]
        }
        places_payload = {"elements": []}
        with patch("osm_fetcher._run_overpass_query", side_effect=[streets_payload, places_payload]):
            result = fetch_osm_data(polygon)
        self.assertTrue(result["streets_ok"])
        self.assertTrue(result["places_ok"])
        self.assertEqual(len(result["streets"]), 1)
        geom = result["streets"][0]["geometry"]
        self.assertAlmostEqual(geom[0][1], 0.0, places=7)
        self.assertAlmostEqual(geom[-1][1], 10.0, places=7)
        self.assertTrue(result["streets"][0]["is_unnamed"])

        with patch("osm_fetcher._run_overpass_query", side_effect=[streets_payload, RuntimeError("places failed")]):
            result = fetch_osm_data(polygon)
        self.assertTrue(result["streets_ok"])
        self.assertFalse(result["places_ok"])
        self.assertEqual(len(result["streets"]), 1)


if __name__ == "__main__":
    unittest.main()
