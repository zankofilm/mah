import os
import tempfile
import unittest

import openpyxl
from pptx import Presentation

from database import Database
from report_generator import (
    generate_block_full_report_pdf,
    generate_block_full_report_excel,
    generate_block_full_report_pptx,
)
from block_report_preview import build_block_full_report_preview_html
from report_preview_html import build_zone_full_report_preview_html


class CommitteeReportTests(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.db = Database(os.path.join(self.tmp.name, "app.db"))
        self.zone_id = self.db.create_zone(
            "بلوک گزارش کمیته",
            [(34.0, 46.0), (34.0, 46.01), (34.01, 46.01), (34.01, 46.0)],
        )
        committee = self.db.get_zone_committees(self.zone_id)[0]
        self.committee_id = committee["id"]
        self.db.update_committee(
            self.committee_id,
            chair_name="علی رئیس کمیته",
            secretary_name="مریم دبیر کمیته",
            status="فعال",
        )
        self.db.add_committee_member(
            self.committee_id,
            "علی رئیس کمیته",
            national_code="0013547899",
            mobile="09120000000",
            member_role="رئیس",
            member_type="نماینده دستگاه",
            agency_name="شهرداری جوانرود",
            is_chair=True,
            status="فعال",
        )
        self.db.add_committee_member(
            self.committee_id,
            "مریم دبیر کمیته",
            national_code="0023547894",
            mobile="09121111111",
            member_role="دبیر",
            member_type="عضو مردمی",
            agency_name="شورای محله",
            is_secretary=True,
            status="فعال",
        )
        meeting_id = self.db.add_committee_meeting(
            self.committee_id, self.zone_id, "جلسه آزمایشی کمیته", meeting_date="2026-07-20"
        )
        self.db.add_committee_resolution(
            self.committee_id,
            self.zone_id,
            "مصوبه آزمایشی کمیته",
            meeting_id=meeting_id,
            responsible_person="علی رئیس کمیته",
            responsible_agency="شهرداری جوانرود",
            due_date="2026-08-01",
        )

    def tearDown(self):
        self.db.close()
        self.tmp.cleanup()

    def test_committee_data_in_all_block_report_formats(self):
        pdf_path = os.path.join(self.tmp.name, "block.pdf")
        xlsx_path = os.path.join(self.tmp.name, "block.xlsx")
        pptx_path = os.path.join(self.tmp.name, "block.pptx")

        generate_block_full_report_pdf(self.db, self.zone_id, pdf_path)
        generate_block_full_report_excel(self.db, self.zone_id, xlsx_path)
        generate_block_full_report_pptx(self.db, self.zone_id, pptx_path)

        self.assertGreater(os.path.getsize(pdf_path), 1000)

        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        self.assertIn("کمیته‌های شش‌گانه", wb.sheetnames)
        self.assertIn("اعضای کمیته‌ها", wb.sheetnames)
        member_values = [cell.value for row in wb["اعضای کمیته‌ها"].iter_rows() for cell in row]
        self.assertIn("علی رئیس کمیته", member_values)
        self.assertTrue(any("رئیس" in str(v) for v in member_values if v))
        self.assertIn("شهرداری جوانرود", member_values)

        prs = Presentation(pptx_path)
        ppt_text = "\n".join(
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        )
        self.assertIn("علی رئیس کمیته", ppt_text)
        self.assertIn("کمیته", ppt_text)

        block_html = build_block_full_report_preview_html(self.db, self.zone_id)
        zone_html = build_zone_full_report_preview_html(self.db, self.zone_id)
        for html in (block_html, zone_html):
            self.assertIn("کمیته‌های شش‌گانه", html)
            self.assertIn("علی رئیس کمیته", html)
            self.assertIn("شهرداری جوانرود", html)
            self.assertIn("مصوبه آزمایشی کمیته", html)


if __name__ == "__main__":
    unittest.main()
