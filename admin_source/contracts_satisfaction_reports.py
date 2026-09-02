# -*- coding: utf-8 -*-
"""گزارش‌های قرارداد، پیمانکار، پرداخت، رضایت و مشارکت مردمی — نسخه ۶.۸."""
from jalali_utils import convert_dates_in_text, format_jalali, now_jalali, install_openpyxl_jalali_patch, install_pptx_jalali_patch
install_openpyxl_jalali_patch()
install_pptx_jalali_patch()
from datetime import datetime


def _safe(value):
    return "—" if value in (None, "") else value


def export_contract_management_excel(db, path, zone_id=None):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    summary = db.get_contract_management_summary(zone_id)
    contractors = db.get_contractors()
    contracts = db.get_contracts(zone_id=zone_id)
    contract_ids = {x["id"] for x in contracts}
    payments = [x for x in db.get_contract_payments() if x["contract_id"] in contract_ids]
    evaluations = [x for x in db.get_contractor_evaluations() if x["contract_id"] in contract_ids]
    surveys = db.get_satisfaction_surveys(zone_id=zone_id)
    participations = db.get_community_participations(zone_id=zone_id)
    alerts = db.get_contract_management_alerts(zone_id=zone_id)

    wb = Workbook(); wb.remove(wb.active)
    navy, gold, border = "13294B", "C9A227", "D7DBE3"

    def style(ws, widths):
        ws.sheet_view.rightToLeft = True
        ws.freeze_panes = "A2"
        for cell in ws[1]:
            cell.fill = PatternFill("solid", fgColor=navy)
            cell.font = Font(color="FFFFFF", bold=True)
            cell.alignment = Alignment(horizontal="center", vertical="center")
        for row in ws.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(horizontal="right", vertical="top", wrap_text=True)
                cell.border = Border(bottom=Side(style="hair", color=border))
        for i, width in enumerate(widths, 1): ws.column_dimensions[chr(64+i)].width = width

    ws = wb.create_sheet("خلاصه")
    ws.append(["شاخص", "مقدار"])
    labels = [
        ("تعداد پیمانکار", summary["contractors_count"]), ("تعداد قرارداد", summary["contracts_count"]),
        ("قرارداد فعال", summary["active_contracts"]), ("مبلغ قراردادها", summary["contract_amount"]),
        ("مبلغ تأییدشده", summary["approved_payments"]), ("مبلغ پرداخت‌شده", summary["paid_amount"]),
        ("مانده قرارداد", summary["remaining_amount"]), ("درصد پرداخت", summary["payment_percent"]),
        ("میانگین رضایت", summary["average_satisfaction"]), ("پاسخ‌دهندگان", summary["survey_respondents"]),
        ("مشارکت‌های مردمی", summary["participations_count"]), ("داوطلبان", summary["volunteers_count"]),
        ("ارزش مشارکت", summary["community_value"]), ("هشدار", summary["alerts_count"]),
    ]
    for row in labels: ws.append(row)
    style(ws, [32, 24])

    ws = wb.create_sheet("پیمانکاران")
    ws.append(["نام", "شناسه ملی", "مدیر", "تلفن", "تخصص", "وضعیت", "امتیاز", "قرارداد", "فعال", "یادداشت"])
    for x in contractors: ws.append([x.get("name"),x.get("national_id"),x.get("manager_name"),x.get("phone"),x.get("specialty"),x.get("status"),x.get("average_score"),x.get("contracts_count"),x.get("active_contracts"),x.get("notes")])
    style(ws, [28,18,22,18,26,14,12,12,12,40])

    ws = wb.create_sheet("قراردادها")
    ws.append(["شماره", "عنوان", "پیمانکار", "پروژه/اقدام", "بلوک", "شروع", "پایان", "مبلغ", "پرداخت", "مانده", "درصد پرداخت", "امتیاز", "وضعیت"])
    for x in contracts: ws.append([x.get("contract_no"),x.get("title"),x.get("contractor_name"),x.get("project_title") or x.get("action_title"),x.get("zone_name"),x.get("start_date"),x.get("end_date"),x.get("amount"),x.get("paid_total"),x.get("remaining_amount"),x.get("payment_percent"),x.get("evaluation_score"),x.get("status")])
    style(ws, [18,34,24,34,18,14,14,18,18,18,14,12,16])

    ws = wb.create_sheet("صورت وضعیت و پرداخت")
    ws.append(["قرارداد", "نوع", "شماره", "دوره", "ناخالص", "کسور", "خالص", "تأیید", "پرداخت", "تاریخ پرداخت", "وضعیت", "یادداشت"])
    for x in payments: ws.append([x.get("contract_no"),x.get("payment_type"),x.get("statement_no"),f"{_safe(x.get('period_from'))} تا {_safe(x.get('period_to'))}",x.get("gross_amount"),x.get("deductions"),x.get("net_amount"),x.get("approved_amount"),x.get("paid_amount"),x.get("payment_date"),x.get("status"),x.get("notes")])
    style(ws, [18,18,14,28,16,16,16,16,16,14,16,40])

    ws = wb.create_sheet("ارزیابی پیمانکار")
    ws.append(["قرارداد", "پیمانکار", "تاریخ", "کیفیت", "زمان", "ایمنی", "همکاری", "مستندات", "امتیاز کل", "ارزیاب", "یادداشت"])
    for x in evaluations: ws.append([x.get("contract_no"),x.get("contractor_name"),x.get("evaluation_date"),x.get("quality_score"),x.get("schedule_score"),x.get("safety_score"),x.get("cooperation_score"),x.get("documentation_score"),x.get("total_score"),x.get("evaluator"),x.get("notes")])
    style(ws, [18,24,14,10,10,10,10,10,12,20,40])

    ws = wb.create_sheet("رضایت مردمی")
    ws.append(["بلوک", "پروژه/اقدام", "تاریخ", "پاسخ‌دهنده", "رفع مسئله٪", "کیفیت", "سرعت", "پاسخ‌گویی", "کلی", "رضایت٪", "بازگشایی", "ثبت‌کننده", "توضیحات"])
    for x in surveys: ws.append([x.get("zone_name"),x.get("project_title") or x.get("action_title") or x.get("tracking_code"),x.get("survey_date"),x.get("respondents"),x.get("problem_resolved_percent"),x.get("quality_score"),x.get("speed_score"),x.get("communication_score"),x.get("overall_score"),x.get("satisfaction_percent"),"بله" if x.get("reopen_recommended") else "خیر",x.get("recorded_by"),x.get("comments")])
    style(ws, [18,34,14,12,14,10,10,12,10,12,12,20,44])

    ws = wb.create_sheet("مشارکت مردمی")
    ws.append(["بلوک", "عنوان", "نوع", "گروه/سازمان", "رابط", "داوطلب", "نقدی", "غیرنقدی", "شروع", "پایان", "وضعیت", "توضیحات"])
    for x in participations: ws.append([x.get("zone_name"),x.get("title"),x.get("participation_type"),x.get("organization_name"),x.get("contact_person"),x.get("volunteers_count"),x.get("cash_value"),x.get("noncash_value"),x.get("start_date"),x.get("end_date"),x.get("status"),x.get("description")])
    style(ws, [18,32,18,28,20,12,16,16,14,14,16,44])

    ws = wb.create_sheet("هشدارها")
    ws.append(["شدت", "نوع", "عنوان", "بلوک", "سررسید", "پیام"])
    for x in alerts: ws.append([x.get("severity"),x.get("type"),x.get("title"),x.get("zone_name"),x.get("due_date"),x.get("message")])
    style(ws, [12,24,36,18,14,60])
    wb.save(path); return path


def export_contract_management_pdf(db, path, zone_id=None):
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    import report_generator as rg
    from report_generator import _pdf_header_flowables, _pdf_styles
    from pdf_text_utils import shape_fa

    summary=db.get_contract_management_summary(zone_id);contracts=db.get_contracts(zone_id=zone_id)
    surveys=db.get_satisfaction_surveys(zone_id=zone_id);parts=db.get_community_participations(zone_id=zone_id)
    alerts=db.get_contract_management_alerts(zone_id=zone_id)
    rg._ensure_fonts_registered()
    cell=ParagraphStyle("c68",fontName=rg.FONT_NAME,fontSize=7.3,leading=9.5,alignment=2,wordWrap="RTL")
    head=ParagraphStyle("h68",fontName=rg.FONT_NAME_BOLD,fontSize=8,leading=10,alignment=2,textColor=colors.white,wordWrap="RTL")
    def table(headers,rows,widths):
        data=[[Paragraph(shape_fa(str(x)),head) for x in headers]]
        data += [[Paragraph(shape_fa(str(_safe(x))),cell) for x in row] for row in rows]
        t=Table(data,colWidths=widths,repeatRows=1,splitByRow=1)
        t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#13294b")),("GRID",(0,0),(-1,-1),.35,colors.HexColor("#d7dbe3")),("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f2f4f7")]),("VALIGN",(0,0),(-1,-1),"MIDDLE"),("ALIGN",(0,0),(-1,-1),"RIGHT"),("PADDING",(0,0),(-1,-1),3)]));return t
    doc=SimpleDocTemplate(path,pagesize=landscape(A4),rightMargin=10*mm,leftMargin=10*mm,topMargin=10*mm,bottomMargin=10*mm)
    _,_,section,normal=_pdf_styles();story=_pdf_header_flowables("گزارش قراردادها، پیمانکاران و رضایت مردمی")
    metrics=[["قرارداد",summary["contracts_count"],"فعال",summary["active_contracts"],"پیمانکار",summary["contractors_count"]],["مبلغ",f"{summary['contract_amount']:,.0f}","پرداخت",f"{summary['paid_amount']:,.0f}","مانده",f"{summary['remaining_amount']:,.0f}"],["رضایت",f"{summary['average_satisfaction']:.1f}٪","پاسخ‌دهنده",summary["survey_respondents"],"داوطلب",summary["volunteers_count"]]]
    story.append(table(["شاخص","مقدار","شاخص","مقدار","شاخص","مقدار"],metrics,[32*mm,25*mm]*3));story.append(Spacer(1,4*mm))
    story.append(Paragraph(shape_fa("قراردادها"),section))
    rows=[[x.get("contract_no"),x.get("title"),x.get("contractor_name"),x.get("zone_name"),x.get("end_date"),f"{float(x.get('amount') or 0):,.0f}",f"{float(x.get('paid_total') or 0):,.0f}",x.get("status")] for x in contracts]
    if rows: story.append(table(["شماره","عنوان","پیمانکار","بلوک","پایان","مبلغ","پرداخت","وضعیت"],rows,[25*mm,52*mm,38*mm,28*mm,24*mm,32*mm,32*mm,28*mm]))
    story.append(PageBreak());story.extend(_pdf_header_flowables("رضایت و مشارکت مردمی"));story.append(Paragraph(shape_fa("رضایت مردمی"),section))
    rows=[[x.get("zone_name"),x.get("project_title") or x.get("action_title"),x.get("survey_date"),x.get("respondents"),f"{float(x.get('satisfaction_percent') or 0):.0f}٪","بله" if x.get("reopen_recommended") else "خیر",x.get("comments")] for x in surveys]
    if rows:story.append(table(["بلوک","پروژه/اقدام","تاریخ","نفر","رضایت","بازگشایی","توضیحات"],rows,[28*mm,60*mm,24*mm,15*mm,20*mm,22*mm,92*mm]))
    story.append(Spacer(1,4*mm));story.append(Paragraph(shape_fa("مشارکت مردمی"),section))
    rows=[[x.get("zone_name"),x.get("title"),x.get("participation_type"),x.get("organization_name"),x.get("volunteers_count"),f"{float(x.get('cash_value') or 0)+float(x.get('noncash_value') or 0):,.0f}",x.get("status")] for x in parts]
    if rows:story.append(table(["بلوک","عنوان","نوع","گروه","داوطلب","ارزش","وضعیت"],rows,[28*mm,60*mm,35*mm,50*mm,20*mm,35*mm,28*mm]))
    if alerts:
        story.append(Spacer(1,4*mm));story.append(Paragraph(shape_fa("هشدارها"),section));story.append(table(["شدت","نوع","عنوان","بلوک","پیام"],[[x.get("severity"),x.get("type"),x.get("title"),x.get("zone_name"),x.get("message")] for x in alerts],[22*mm,38*mm,55*mm,30*mm,100*mm]))
    doc.build(story);return path


def export_contract_management_powerpoint(db, path, zone_id=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    summary=db.get_contract_management_summary(zone_id);contracts=db.get_contracts(zone_id=zone_id)
    surveys=db.get_satisfaction_surveys(zone_id=zone_id);parts=db.get_community_participations(zone_id=zone_id)
    prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5)
    def add_slide(title, lines):
        slide=prs.slides.add_slide(prs.slide_layouts[5]);tb=slide.shapes.add_textbox(Inches(.6),Inches(.35),Inches(12.1),Inches(.7));p=tb.text_frame.paragraphs[0];p.text=title;p.font.size=Pt(28);p.font.bold=True;p.alignment=PP_ALIGN.RIGHT
        body=slide.shapes.add_textbox(Inches(.8),Inches(1.2),Inches(11.7),Inches(5.7));tf=body.text_frame;tf.word_wrap=True
        for i,line in enumerate(lines):
            par=tf.paragraphs[0] if i==0 else tf.add_paragraph();par.text=str(line);par.font.size=Pt(18);par.alignment=PP_ALIGN.RIGHT;par.space_after=Pt(9)
    add_slide("قراردادها، پیمانکاران و رضایت مردمی",[f"تاریخ تولید: {format_jalali(datetime.now().strftime('%Y-%m-%d'))}",f"قراردادها: {summary['contracts_count']} | فعال: {summary['active_contracts']}",f"مبلغ قرارداد: {summary['contract_amount']:,.0f}",f"پرداخت: {summary['paid_amount']:,.0f} | مانده: {summary['remaining_amount']:,.0f}",f"میانگین رضایت: {summary['average_satisfaction']:.1f}٪"])
    add_slide("قراردادهای اصلی",[f"{x.get('contract_no')} — {x.get('title')} | {x.get('contractor_name')} | پرداخت {float(x.get('payment_percent') or 0):.0f}٪ | {x.get('status')}" for x in contracts[:12]] or ["قراردادی ثبت نشده است."])
    add_slide("رضایت مردم",[f"{x.get('zone_name')} — {x.get('project_title') or x.get('action_title') or 'خدمت'} | {float(x.get('satisfaction_percent') or 0):.0f}٪ | {x.get('respondents')} نفر" for x in surveys[:12]] or ["نظرسنجی ثبت نشده است."])
    add_slide("مشارکت مردمی",[f"{x.get('zone_name')} — {x.get('title')} | {x.get('volunteers_count')} داوطلب | ارزش {float(x.get('cash_value') or 0)+float(x.get('noncash_value') or 0):,.0f}" for x in parts[:12]] or ["مشارکتی ثبت نشده است."])
    prs.save(path);return path


def build_contract_management_preview_html(db, zone_id=None):
    import html
    s=db.get_contract_management_summary(zone_id);contracts=db.get_contracts(zone_id=zone_id);surveys=db.get_satisfaction_surveys(zone_id=zone_id);parts=db.get_community_participations(zone_id=zone_id);alerts=db.get_contract_management_alerts(zone_id=zone_id)
    def esc(v):return html.escape(convert_dates_in_text(str(_safe(v))))
    def rows(data, keys):
        return ''.join('<tr>'+''.join(f'<td>{esc(k(x) if callable(k) else x.get(k))}</td>' for k in keys)+'</tr>' for x in data) or '<tr><td colspan="10">رکوردی ثبت نشده است.</td></tr>'
    return f'''<!doctype html><html lang="fa" dir="rtl"><head><meta charset="utf-8"><style>
    body{{font-family:Tahoma;background:#eef1f5;color:#1c2530;margin:0;padding:24px}}.page{{max-width:1200px;margin:auto;background:white;padding:28px;border-radius:12px}}
    h1,h2{{color:#13294b}}.kpi{{display:grid;grid-template-columns:repeat(4,1fr);gap:10px}}.card{{border:1px solid #d7dbe3;border-radius:9px;padding:14px;background:linear-gradient(#fff,#edf1f6)}}
    table{{width:100%;border-collapse:collapse;margin:12px 0 24px}}th{{background:#13294b;color:white}}th,td{{border:1px solid #d7dbe3;padding:7px;text-align:right}}tr:nth-child(even){{background:#f4f6f8}}
    </style></head><body><div class="page"><h1>گزارش قراردادها، پیمانکاران و رضایت مردمی</h1>
    <div class="kpi"><div class="card">قراردادها<br><b>{s['contracts_count']}</b></div><div class="card">مبلغ قرارداد<br><b>{s['contract_amount']:,.0f}</b></div><div class="card">پرداخت<br><b>{s['payment_percent']:.0f}٪</b></div><div class="card">رضایت مردم<br><b>{s['average_satisfaction']:.0f}٪</b></div></div>
    <h2>قراردادها</h2><table><tr><th>شماره</th><th>عنوان</th><th>پیمانکار</th><th>بلوک</th><th>مبلغ</th><th>پرداخت</th><th>وضعیت</th></tr>{rows(contracts,['contract_no','title','contractor_name','zone_name',lambda x:f"{float(x.get('amount') or 0):,.0f}",lambda x:f"{float(x.get('payment_percent') or 0):.0f}٪",'status'])}</table>
    <h2>رضایت مردمی</h2><table><tr><th>بلوک</th><th>پروژه/اقدام</th><th>پاسخ‌دهنده</th><th>رضایت</th><th>بازگشایی</th><th>توضیحات</th></tr>{rows(surveys,['zone_name',lambda x:x.get('project_title') or x.get('action_title'),'respondents',lambda x:f"{float(x.get('satisfaction_percent') or 0):.0f}٪",lambda x:'بله' if x.get('reopen_recommended') else 'خیر','comments'])}</table>
    <h2>مشارکت مردمی</h2><table><tr><th>بلوک</th><th>عنوان</th><th>نوع</th><th>گروه</th><th>داوطلب</th><th>ارزش</th></tr>{rows(parts,['zone_name','title','participation_type','organization_name','volunteers_count',lambda x:f"{float(x.get('cash_value') or 0)+float(x.get('noncash_value') or 0):,.0f}"])}</table>
    <h2>هشدارها</h2><table><tr><th>شدت</th><th>نوع</th><th>عنوان</th><th>پیام</th></tr>{rows(alerts,['severity','type','title','message'])}</table>
    </div></body></html>'''
