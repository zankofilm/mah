# -*- coding: utf-8 -*-
"""
ماژول تولید گزارش‌های PDF و Excel بر اساس داده‌های دیتابیس.

گزارش‌های طراحی‌شده:
  1. گزارش کلی وضعیت (آمار خلاصه همه بخش‌ها)
  2. گزارش کامل یک منطقه/بلوک (خیابان‌ها، اماکن، اعضای شورا، محل جلسات، درخواست‌ها)
  3. گزارش اعضای شورای محلات (همه مناطق یا یک منطقه خاص)
  4. گزارش درخواست‌ها و مشکلات اولویت‌بندی‌شده (با وضعیت پیگیری)
  5. گزارش اقدامات انجام‌شده (تاریخچه کامل پیگیری هر درخواست)
"""

from jalali_utils import convert_dates_in_text, format_jalali, now_jalali, install_openpyxl_jalali_patch, install_pptx_jalali_patch
import zone_action_plan
install_openpyxl_jalali_patch()
install_pptx_jalali_patch()
import os
from runtime_paths import get_temp_dir
import io
from datetime import datetime

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.lib.units import mm
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image, KeepTogether
)
from reportlab.lib.styles import ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill
from openpyxl.utils import get_column_letter

from pdf_text_utils import shape_fa
from committee_report_utils import (
    get_zone_committee_report_data, committee_summary_rows, committee_member_rows,
    committee_meeting_rows, committee_resolution_rows, committee_link_rows, member_display_role,
)


def _snapshot_image_source(db, zone_id, map_image_path=None):
    """مسیر فایل یا BytesIO تصویر ذخیره‌شده بلوک را برای گزارش برمی‌گرداند."""
    snapshot = None
    try:
        from zone_snapshot_service import refresh_zone_snapshot
        snapshot = refresh_zone_snapshot(db, zone_id, force=False)
    except Exception:
        snapshot = db.get_zone_snapshot(zone_id) if hasattr(db, "get_zone_snapshot") else None
    if snapshot and snapshot.get("png_data"):
        return io.BytesIO(snapshot["png_data"]), snapshot
    if map_image_path and os.path.exists(map_image_path):
        return map_image_path, None
    return None, snapshot

FONT_DIR = os.path.join(os.path.dirname(__file__), "fonts")
FONT_REGULAR_PATH = os.path.join(FONT_DIR, "Vazirmatn-Regular.ttf")
FONT_BOLD_PATH = os.path.join(FONT_DIR, "Vazirmatn-Bold.ttf")

# فونت در بسته برنامه قرار داده نمی‌شود؛ ابتدا فونت اختیاری کاربر و سپس فونت‌های
# استاندارد سیستم‌عامل بررسی می‌شوند تا متن فارسی به مربع تبدیل نشود.
_FONT_CANDIDATES = [
    (FONT_REGULAR_PATH, FONT_BOLD_PATH, "Vazirmatn"),
    (r"C:\\Windows\\Fonts\\tahoma.ttf", r"C:\\Windows\\Fonts\\tahomabd.ttf", "TahomaFa"),
    (r"C:\\Windows\\Fonts\\arial.ttf", r"C:\\Windows\\Fonts\\arialbd.ttf", "ArialFa"),
    (r"C:\\Windows\\Fonts\\segoeui.ttf", r"C:\\Windows\\Fonts\\segoeuib.ttf", "SegoeUIFa"),
    ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", "DejaVuFa"),
    ("/usr/share/fonts/truetype/noto/NotoSansArabic-Regular.ttf", "/usr/share/fonts/truetype/noto/NotoSansArabic-Bold.ttf", "NotoArabicFa"),
    ("/System/Library/Fonts/Supplemental/Arial.ttf", "/System/Library/Fonts/Supplemental/Arial Bold.ttf", "MacArialFa"),
]

_FONT_REGISTERED = False
_FONT_SOURCE_PATH = None
FONT_NAME = "Helvetica"
FONT_NAME_BOLD = "Helvetica-Bold"


def _ensure_fonts_registered():
    """ثبت نخستین فونت فارسی قابل استفاده از پوشه برنامه یا فونت‌های سیستم."""
    global _FONT_REGISTERED, _FONT_SOURCE_PATH, FONT_NAME, FONT_NAME_BOLD
    if _FONT_REGISTERED:
        return
    for regular_path, bold_path, base_name in _FONT_CANDIDATES:
        if not os.path.exists(regular_path):
            continue
        try:
            regular_name = base_name + "-Regular"
            bold_name = base_name + "-Bold"
            pdfmetrics.registerFont(TTFont(regular_name, regular_path))
            FONT_NAME = regular_name
            if os.path.exists(bold_path):
                pdfmetrics.registerFont(TTFont(bold_name, bold_path))
                FONT_NAME_BOLD = bold_name
            else:
                FONT_NAME_BOLD = regular_name
            _FONT_SOURCE_PATH = regular_path
            break
        except Exception:
            continue
    _FONT_REGISTERED = True


def fonts_missing():
    """True فقط وقتی هیچ فونت فارسی محلی یا سیستمی قابل ثبت نباشد."""
    _ensure_fonts_registered()
    return _FONT_SOURCE_PATH is None


# ==================== PDF ====================

def _pdf_styles():
    _ensure_fonts_registered()
    title_style = ParagraphStyle(
        "TitleFa", fontName=FONT_NAME_BOLD, fontSize=15,
        alignment=1, leading=22, textColor=colors.HexColor("#0b1f3a")
    )
    subtitle_style = ParagraphStyle(
        "SubtitleFa", fontName=FONT_NAME, fontSize=10,
        alignment=1, leading=16, textColor=colors.HexColor("#5b6472")
    )
    section_style = ParagraphStyle(
        "SectionFa", fontName=FONT_NAME_BOLD, fontSize=12,
        alignment=2, leading=18, textColor=colors.HexColor("#13294b"),
        spaceBefore=10, spaceAfter=6
    )
    normal_style = ParagraphStyle(
        "NormalFa", fontName=FONT_NAME, fontSize=9,
        alignment=2, leading=14
    )
    return title_style, subtitle_style, section_style, normal_style


def _pdf_header_flowables(report_title):
    title_style, subtitle_style, section_style, normal_style = _pdf_styles()
    flowables = []
    flowables.append(Paragraph(shape_fa("فرمانداری شهرستان جوانرود"), title_style))
    flowables.append(Paragraph(shape_fa("وزارت کشور — استانداری کرمانشاه"), subtitle_style))
    flowables.append(Spacer(1, 6 * mm))
    flowables.append(Paragraph(shape_fa(report_title), section_style))
    now_str = now_jalali()
    flowables.append(Paragraph(shape_fa(f"تاریخ تهیه گزارش: {now_str}"), normal_style))
    flowables.append(Spacer(1, 4 * mm))
    return flowables


def _make_pdf_table(headers, rows, col_widths=None):
    """جدول PDF با راست‌چین کردن متن فارسی هر سلول."""
    _ensure_fonts_registered()
    shaped_headers = [shape_fa(h) for h in headers]
    shaped_rows = [[shape_fa(cell) for cell in row] for row in rows]
    data = [shaped_headers] + shaped_rows

    table = Table(data, colWidths=col_widths, repeatRows=1)
    table.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), FONT_NAME),
        ("FONTNAME", (0, 0), (-1, 0), FONT_NAME_BOLD),
        ("FONTSIZE", (0, 0), (-1, -1), 8.5),
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#13294b")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("ALIGN", (0, 0), (-1, -1), "RIGHT"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d7dbe3")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f0f2f6")]),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    return table


def generate_overview_report_pdf(db, output_path):
    """گزارش کلی وضعیت سامانه (آمار خلاصه همه بخش‌ها)."""
    doc = SimpleDocTemplate(output_path, pagesize=A4,
                             topMargin=15 * mm, bottomMargin=15 * mm,
                             leftMargin=15 * mm, rightMargin=15 * mm)
    flowables = _pdf_header_flowables("گزارش کلی وضعیت سامانه")

    zones = db.get_zones()
    all_streets = db.get_streets()
    all_places = db.get_places()
    all_members = db.get_council_members()
    all_requests = db.get_priority_requests()
    all_committees = [c for z in zones for c in get_zone_committee_report_data(db, z["id"])]

    total_actions = sum(r["action_count"] for r in all_requests)
    resolved_like = [r for r in all_requests if r["action_count"] > 0]

    headers = ["شاخص", "مقدار"]
    rows = [
        ["تعداد مناطق/بلوک‌های ثبت‌شده", str(len(zones))],
        ["تعداد خیابان‌ها و کوچه‌های ثبت‌شده", str(len(all_streets))],
        ["تعداد اماکن ثبت‌شده (مدارس، ادارات و ...)", str(len(all_places))],
        ["تعداد مساجد مرجع", str(len(db.get_mosques()))],
        ["تعداد اعضای شورای محلات", str(len(all_members))],
        ["تعداد کمیته‌های تخصصی", str(len(all_committees))],
        ["تعداد اعضای فعال کمیته‌ها", str(sum(len(c.get("active_members") or []) for c in all_committees))],
        ["تعداد مصوبات باز کمیته‌ها", str(sum(1 for c in all_committees for r in (c.get("resolutions") or []) if r.get("status") not in ("انجام‌شده", "لغوشده")))],
        ["تعداد درخواست‌ها و مشکلات ثبت‌شده", str(len(all_requests))],
        ["تعداد درخواست‌های دارای حداقل یک اقدام پیگیری", str(len(resolved_like))],
        ["مجموع اقدامات پیگیری ثبت‌شده", str(total_actions)],
    ]
    flowables.append(_make_pdf_table(headers, rows, col_widths=[110 * mm, 60 * mm]))
    flowables.append(Spacer(1, 8 * mm))

    if zones:
        _, _, section_style, _ = _pdf_styles()
        flowables.append(Paragraph(shape_fa("جزئیات هر منطقه"), section_style))
        zone_headers = ["نام منطقه", "مساحت (هکتار)", "خیابان", "مکان", "مسجد", "درخواست"]
        zone_rows = []
        for z in zones:
            street_count = len(db.get_streets(zone_id=z["id"]))
            place_count = len(db.get_places(zone_id=z["id"]))
            mosque_count = len(db.get_mosques(zone_id=z["id"]))
            request_count = len(db.get_priority_requests(zone_id=z["id"]))
            zone_rows.append([z["name"], f"{(z.get('area_m2', 0) or 0)/10000:.2f}", str(street_count), str(place_count), str(mosque_count), str(request_count)])
        flowables.append(_make_pdf_table(zone_headers, zone_rows,
                                          col_widths=[48 * mm, 28 * mm, 24 * mm, 24 * mm, 22 * mm, 24 * mm]))

    doc.build(flowables)
    return output_path


def generate_block_full_report_pdf(db, zone_id, output_path, map_image_path=None):
    """گزارش کامل بلوک با تصویر گرافیکی در ابتدای گزارش و سپس همه مشخصات به ترتیب."""
    zone = db.get_zone(zone_id)
    if not zone:
        raise ValueError("منطقه یافت نشد")

    doc = SimpleDocTemplate(
        output_path, pagesize=A4,
        topMargin=12 * mm, bottomMargin=15 * mm,
        leftMargin=15 * mm, rightMargin=15 * mm,
    )
    flowables = _pdf_header_flowables(f"گزارش کامل بلوک: {zone['name']}")
    _, _, section_style, normal_style = _pdf_styles()

    # ۱) تصویر گرافیکی بلوک؛ نخستین محتوای گزارش و حداکثر در ابعاد نیم صفحه A4.
    flowables.append(Paragraph(shape_fa("نمای گرافیکی بلوک"), section_style))
    image_source, snapshot_meta = _snapshot_image_source(db, zone_id, map_image_path)
    if image_source is not None:
        try:
            source_width = float((snapshot_meta or {}).get("width", 1200) or 1200)
            source_height = float((snapshot_meta or {}).get("height", 900) or 900)
            max_width = 180 * mm
            max_height = 125 * mm
            scale = min(max_width / source_width, max_height / source_height)
            report_width = source_width * scale
            report_height = source_height * scale
            img = Image(image_source, width=report_width, height=report_height)
            img.hAlign = "CENTER"
            flowables.append(img)
            if snapshot_meta:
                flowables.append(Paragraph(
                    shape_fa(
                        f"نسخه تصویر: {snapshot_meta.get('version', 1)} — "
                        f"تاریخ تولید: {snapshot_meta.get('generated_at') or '—'}"
                    ),
                    normal_style,
                ))
        except Exception:
            flowables.append(Paragraph(shape_fa("خطا در بارگذاری تصویر گرافیکی بلوک."), normal_style))
    else:
        flowables.append(Paragraph(shape_fa("نمای گرافیکی بلوک هنوز تولید نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    streets = db.get_streets(zone_id=zone_id)
    places = db.get_places(zone_id=zone_id)
    mosques = db.get_mosques(zone_id=zone_id)
    members = db.get_council_members(zone_id=zone_id)
    requests = db.get_priority_requests(zone_id=zone_id)
    meeting_place = db.get_zone_meeting_place(zone_id)
    profile = db.get_zone_profile(zone_id)
    neighborhood_issues = db.get_neighborhood_issues(zone_id)
    neighborhood_actions = db.get_neighborhood_actions(zone_id)
    neighborhood_meetings = db.get_neighborhood_meetings(zone_id)
    neighborhood_resolutions = db.get_neighborhood_resolutions(zone_id=zone_id)
    neighborhood_budgets = db.get_neighborhood_budgets(zone_id)
    budget_summary = db.get_budget_summary(zone_id)
    performance = db.get_zone_performance(zone_id)
    management_alerts = db.get_management_alerts(zone_id)
    quality_issues = db.get_quality_issues(zone_id)
    field_visits = db.get_field_visits(zone_id)
    citizen_requests = db.get_citizen_requests(zone_id)
    operational_analysis = db.get_zone_operational_analysis(zone_id)
    correspondence_letters = db.get_correspondence_letters(zone_id=zone_id)
    approval_requests = db.get_approval_requests(zone_id=zone_id, limit=5000)
    generated_documents = db.get_generated_documents(zone_id=zone_id, limit=5000)
    committees = get_zone_committee_report_data(db, zone_id)

    # ۲) مشخصات پایه بلوک
    flowables.append(Paragraph(shape_fa("مشخصات پایه بلوک"), section_style))
    summary_rows = [
        ["نام بلوک", zone.get("name") or "—"],
        ["وضعیت تکمیل", zone.get("status") or "—"],
        ["مساحت", f"{(zone.get('area_m2', 0) or 0)/10000:.2f} هکتار"],
        ["محیط", f"{zone.get('perimeter_m', 0) or 0:.0f} متر"],
        ["تعداد نقاط مرزی", str(len(zone.get("boundary_points", [])))],
        ["تعداد خیابان و کوچه", str(len(streets))],
        ["تعداد سایر اماکن", str(len(places))],
        ["تعداد مساجد", str(len(mosques))],
        ["تعداد کمیته‌های تخصصی", str(len(committees))],
        ["اعضای فعال کمیته‌ها", str(sum(len(c.get("active_members") or []) for c in committees))],
        ["مصوبات باز کمیته‌ها", str(sum(1 for c in committees for r in (c.get("resolutions") or []) if r.get("status") not in ("انجام‌شده", "لغوشده")))],
        ["بازدیدهای میدانی", str(len(field_visits))],
        ["درخواست‌های مردمی", str(len(citizen_requests))],
        ["مکاتبات اداری مرتبط", str(len(correspondence_letters))],
        ["گردش‌های تأیید", str(len(approval_requests))],
        ["اسناد Word تولیدشده", str(len(generated_documents))],
        ["سطح ریسک عملیاتی", operational_analysis.get("risk_level") or "—"],
        ["تاریخ ایجاد", zone.get("created_at") or "—"],
        ["آخرین بروزرسانی", zone.get("updated_at") or "—"],
    ]
    flowables.append(_make_pdf_table(["شاخص", "مقدار"], summary_rows, col_widths=[60 * mm, 110 * mm]))
    flowables.append(Spacer(1, 6 * mm))

    # ۲.۱) خلاصه برنامه عملیاتی هوشمند — از آخرین برنامه ذخیره‌شده استفاده می‌شود
    # تا تولید گزارش کند نشود؛ در صورت نبود برنامه ذخیره‌شده، بلادرنگ و آفلاین
    # (بدون فراخوانی شبکه) یک نسخه تازه ساخته می‌شود.
    flowables.append(Paragraph(shape_fa("خلاصه برنامه عملیاتی (تولیدشده توسط سامانه)"), section_style))
    try:
        latest_plan = db.get_latest_zone_action_plan(zone_id)
        if latest_plan:
            plan_text = latest_plan["content"]
            plan_note = f"تولید‌شده در تاریخ {latest_plan.get('created_at') or '—'}"
        else:
            action_plan_context = db.get_zone_action_plan_context(zone_id)
            plan_text = zone_action_plan.generate_offline(action_plan_context)
            plan_note = "تولید آفلاین در لحظه تهیه این گزارش (هنوز برنامه‌ای برای این بلوک ذخیره نشده بود)"
        for line in plan_text.split("\n"):
            if line.strip():
                flowables.append(Paragraph(shape_fa(line.strip()), normal_style))
        flowables.append(Paragraph(shape_fa(plan_note), normal_style))
    except Exception:
        flowables.append(Paragraph(shape_fa("خلاصه برنامه عملیاتی در دسترس نیست."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    # ۳) جمعیت، خانوار و شاخص‌های اجتماعی
    flowables.append(Paragraph(shape_fa("جمعیت و خانوار"), section_style))
    profile_rows = [
        ["ساختمان مسکونی", str(profile.get("residential_buildings", 0))],
        ["واحد مسکونی", str(profile.get("residential_units", 0))],
        ["واحد دارای سکنه", str(profile.get("occupied_units", 0))],
        ["خانوار محاسبه‌شده", str(profile.get("estimated_households", 0))],
        ["خانوار ثبت میدانی", str(profile.get("field_households", 0))],
        ["خانوار نهایی تأییدشده", str(profile.get("approved_households", 0))],
        ["جمعیت تخمینی", str(profile.get("estimated_population", 0))],
        ["روش برآورد", profile.get("estimation_method") or "—"],
        ["سطح اطمینان", profile.get("confidence_level") or "—"],
    ]
    flowables.append(_make_pdf_table(["شاخص", "مقدار"], profile_rows, col_widths=[70 * mm, 100 * mm]))
    flowables.append(Spacer(1, 6 * mm))

    # ۴) محل برگزاری جلسات
    flowables.append(Paragraph(shape_fa("محل برگزاری جلسات شورا"), section_style))
    if meeting_place:
        rows = [
            ["نام مکان", meeting_place["place_name"] or "—"],
            ["آدرس دقیق", meeting_place["exact_address"] or "—"],
        ]
        flowables.append(_make_pdf_table(["فیلد", "مقدار"], rows, col_widths=[50 * mm, 120 * mm]))
    else:
        flowables.append(Paragraph(shape_fa("هنوز محل جلسه‌ای برای این بلوک ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    # ۴) خیابان‌ها و کوچه‌ها
    flowables.append(Paragraph(shape_fa(f"خیابان‌ها و کوچه‌ها (تعداد: {len(streets)})"), section_style))
    if streets:
        street_rows = []
        for index, street in enumerate(streets, start=1):
            kind = street.get("highway_type") or "—"
            street_rows.append([str(index), street.get("name") or "معبر بدون نام", kind])
        flowables.append(_make_pdf_table(
            ["ردیف", "نام معبر", "نوع"], street_rows,
            col_widths=[18 * mm, 105 * mm, 47 * mm],
        ))
    else:
        flowables.append(Paragraph(shape_fa("معبری برای این بلوک ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    # ۵) سایر اماکن داخل بلوک
    flowables.append(Paragraph(shape_fa(f"اماکن داخل بلوک، به‌جز مساجد مرجع (تعداد: {len(places)})"), section_style))
    if places:
        place_rows = [
            [str(i), p.get("name") or "—", p.get("category") or "—", p.get("subtype") or "—"]
            for i, p in enumerate(places, start=1)
        ]
        flowables.append(_make_pdf_table(
            ["ردیف", "نام مکان", "دسته", "نوع"], place_rows,
            col_widths=[16 * mm, 76 * mm, 40 * mm, 38 * mm],
        ))
    else:
        flowables.append(Paragraph(shape_fa("مکان دیگری برای این بلوک ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    # ۶) مساجد داخل بلوک
    flowables.append(Paragraph(shape_fa(f"مساجد داخل بلوک (تعداد: {len(mosques)})"), section_style))
    if mosques:
        mosque_rows = [
            [str(i), m["name"], f"{m['lat']:.6f}", f"{m['lon']:.6f}"]
            for i, m in enumerate(mosques, start=1)
        ]
        flowables.append(_make_pdf_table(
            ["ردیف", "نام مسجد", "عرض", "طول"], mosque_rows,
            col_widths=[16 * mm, 84 * mm, 35 * mm, 35 * mm],
        ))
    else:
        flowables.append(Paragraph(shape_fa("مسجدی داخل محدوده این بلوک ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    # ۷) معتمدین و اعضای شورا
    flowables.append(Paragraph(
        shape_fa(f"نام و نام‌خانوادگی معتمدین و شماره تماس (تعداد: {len(members)})"), section_style
    ))
    if members:
        member_rows = [
            [
                f"{m['first_name']} {m['last_name']}",
                m["member_group"] or "—",
                m["position"] or "—",
                m["mobile"] or "—",
            ]
            for m in members
        ]
        flowables.append(_make_pdf_table(
            ["نام و نام‌خانوادگی", "دسته", "سمت", "شماره تماس"], member_rows,
            col_widths=[55 * mm, 35 * mm, 40 * mm, 40 * mm],
        ))
    else:
        flowables.append(Paragraph(shape_fa("عضوی برای این بلوک ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    # ۸) کمیته‌های شش‌گانه، اعضا و سمت‌ها
    flowables.append(Paragraph(shape_fa(f"کمیته‌های تخصصی بلوک (تعداد: {len(committees)})"), section_style))
    if committees:
        flowables.append(_make_pdf_table(
            ["ردیف", "نام کمیته", "رئیس", "دبیر", "اعضای فعال", "جلسات", "مصوبات باز", "وضعیت"],
            committee_summary_rows(committees),
            col_widths=[9*mm, 42*mm, 28*mm, 28*mm, 17*mm, 14*mm, 18*mm, 17*mm],
        ))
        flowables.append(Spacer(1, 4 * mm))
        for committee in committees:
            flowables.append(Paragraph(shape_fa(committee.get("title") or "کمیته"), section_style))
            committee_profile = [
                ["رئیس کمیته", committee.get("chair_name") or "—"],
                ["دبیر کمیته", committee.get("secretary_name") or "—"],
                ["شماره حکم", committee.get("decree_no") or "—"],
                ["تاریخ حکم", committee.get("decree_date") or "—"],
                ["دوره فعالیت", f"{committee.get('start_date') or '—'} تا {committee.get('end_date') or '—'}"],
                ["دستگاه‌های پیشنهادی", committee.get("recommended_agencies") or "—"],
                ["وضعیت", committee.get("status") or "—"],
            ]
            flowables.append(_make_pdf_table(["مشخصه", "مقدار"], committee_profile, col_widths=[42*mm, 128*mm]))

            committee_members = committee.get("members") or []
            flowables.append(Paragraph(shape_fa(f"اعضا و سمت‌ها (تعداد: {len(committee_members)})"), normal_style))
            if committee_members:
                member_rows = [[
                    member.get("person_name") or "—", member.get("national_code") or "—",
                    member_display_role(member), member.get("member_type") or "—",
                    member.get("agency_name") or "—", member.get("mobile") or "—",
                    member.get("status") or "—",
                ] for member in committee_members]
                flowables.append(_make_pdf_table(
                    ["نام عضو", "کد ملی", "سمت", "نوع عضویت", "اداره/دستگاه", "موبایل", "وضعیت"],
                    member_rows, col_widths=[31*mm, 22*mm, 28*mm, 25*mm, 32*mm, 23*mm, 18*mm],
                ))
            else:
                flowables.append(Paragraph(shape_fa("فاقد عضو ثبت‌شده است."), normal_style))

            committee_meetings = committee.get("meetings") or []
            if committee_meetings:
                meeting_rows = [[
                    m.get("title") or "—", m.get("meeting_date") or "—", m.get("start_time") or "—",
                    m.get("place_name") or "—", m.get("status") or "—",
                ] for m in committee_meetings]
                flowables.append(Paragraph(shape_fa("جلسات کمیته"), normal_style))
                flowables.append(_make_pdf_table(
                    ["عنوان", "تاریخ", "ساعت", "محل", "وضعیت"], meeting_rows,
                    col_widths=[55*mm, 28*mm, 20*mm, 42*mm, 25*mm],
                ))

            committee_resolutions = committee.get("resolutions") or []
            if committee_resolutions:
                resolution_rows = [[
                    r.get("title") or "—", r.get("responsible_person") or "—",
                    r.get("responsible_agency") or "—", r.get("due_date") or "—", r.get("status") or "—",
                ] for r in committee_resolutions]
                flowables.append(Paragraph(shape_fa("مصوبات کمیته"), normal_style))
                flowables.append(_make_pdf_table(
                    ["عنوان مصوبه", "مسئول پیگیری", "دستگاه مسئول", "مهلت", "وضعیت"],
                    resolution_rows, col_widths=[55*mm, 34*mm, 36*mm, 24*mm, 25*mm],
                ))

            linked_rows = []
            linked_rows.extend([["مسئله", x.get("title") or "—", x.get("related_office") or "—", x.get("status") or "—"] for x in committee.get("issues") or []])
            linked_rows.extend([["اقدام", x.get("title") or "—", x.get("responsible_office") or "—", x.get("status") or "—"] for x in committee.get("actions") or []])
            if linked_rows:
                flowables.append(Paragraph(shape_fa("مسائل و اقدامات مرتبط"), normal_style))
                flowables.append(_make_pdf_table(
                    ["نوع", "عنوان", "دستگاه", "وضعیت"], linked_rows,
                    col_widths=[22*mm, 70*mm, 45*mm, 33*mm],
                ))
            flowables.append(Spacer(1, 5 * mm))
    else:
        flowables.append(Paragraph(shape_fa("کمیته‌ای برای این بلوک ایجاد نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    # مکاتبات اداری مرتبط با بلوک
    flowables.append(Paragraph(shape_fa(f"مکاتبات اداری مرتبط (تعداد: {len(correspondence_letters)})"), section_style))
    if correspondence_letters:
        letter_rows = [[
            item.get("letter_number") or "—", item.get("direction") or "—",
            item.get("subject") or "—", item.get("sender") or "—", item.get("recipient") or "—",
            item.get("due_date") or "—", item.get("status") or "—", str(item.get("attachment_count") or 0)
        ] for item in correspondence_letters]
        flowables.append(_make_pdf_table(
            ["شماره", "نوع", "موضوع", "فرستنده", "گیرنده", "مهلت", "وضعیت", "پیوست"],
            letter_rows, col_widths=[22*mm, 18*mm, 42*mm, 24*mm, 24*mm, 20*mm, 22*mm, 12*mm]
        ))
    else:
        flowables.append(Paragraph(shape_fa("مکاتبه‌ای برای این بلوک ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    flowables.append(Paragraph(shape_fa(f"گردش‌های تأیید اداری (تعداد: {len(approval_requests)})"), section_style))
    if approval_requests:
        approval_rows = [[
            item.get("title") or "—", item.get("entity_type") or "—",
            f"{item.get('current_step') or 0}/{item.get('total_steps') or 0}",
            item.get("due_date") or "—", item.get("status") or "—",
            item.get("requested_by_name") or "system"
        ] for item in approval_requests]
        flowables.append(_make_pdf_table(
            ["عنوان", "نوع", "مرحله", "مهلت", "وضعیت", "درخواست‌کننده"],
            approval_rows, col_widths=[58*mm, 25*mm, 18*mm, 25*mm, 28*mm, 30*mm]
        ))
    else:
        flowables.append(Paragraph(shape_fa("گردش تأییدی برای این بلوک ثبت نشده است."), normal_style))
    if generated_documents:
        flowables.append(Paragraph(shape_fa(f"اسناد اداری تولیدشده (تعداد: {len(generated_documents)})"), section_style))
        document_rows = [[
            item.get("title") or "—", item.get("template_name") or "—",
            item.get("created_by_name") or "system", item.get("created_at") or "—"
        ] for item in generated_documents]
        flowables.append(_make_pdf_table(
            ["عنوان سند", "قالب", "تولیدکننده", "تاریخ"], document_rows,
            col_widths=[75*mm, 45*mm, 35*mm, 30*mm]
        ))
    flowables.append(Spacer(1, 6 * mm))

    # عملیات میدانی
    flowables.append(Paragraph(shape_fa(f"بازدیدها و برداشت‌های میدانی (تعداد: {len(field_visits)})"), section_style))
    if field_visits:
        visit_rows = [[
            str(i), item.get("visit_date") or "—", item.get("officer_name") or "—",
            item.get("visit_type") or "—", item.get("location_text") or "—",
            str(item.get("households_count") or 0),
            "بله" if item.get("followup_required") else "خیر", item.get("status") or "—"
        ] for i, item in enumerate(field_visits, start=1)]
        flowables.append(_make_pdf_table(
            ["ردیف", "تاریخ", "کارشناس", "نوع", "موقعیت", "خانوار", "پیگیری", "وضعیت"],
            visit_rows, col_widths=[9*mm, 19*mm, 25*mm, 25*mm, 37*mm, 15*mm, 15*mm, 25*mm]
        ))
    else:
        flowables.append(Paragraph(shape_fa("بازدید میدانی ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    # درخواست‌های مردمی
    flowables.append(Paragraph(shape_fa(f"درخواست‌ها و گزارش‌های مردمی (تعداد: {len(citizen_requests)})"), section_style))
    if citizen_requests:
        citizen_rows = [[
            item.get("tracking_code") or "—", item.get("title") or "—",
            item.get("category") or "—", str(item.get("urgency") or 0),
            item.get("assigned_office") or "—", item.get("status") or "—",
            str(item.get("linked_issue_id") or "—")
        ] for item in citizen_requests]
        flowables.append(_make_pdf_table(
            ["کد رهگیری", "عنوان", "دسته", "فوریت", "دستگاه", "وضعیت", "مسئله"],
            citizen_rows, col_widths=[26*mm, 42*mm, 23*mm, 14*mm, 28*mm, 27*mm, 15*mm]
        ))
    else:
        flowables.append(Paragraph(shape_fa("درخواست مردمی ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    flowables.append(Paragraph(shape_fa("تحلیل عملیاتی بلوک"), section_style))
    analysis_rows = [
        ["امتیاز ریسک", str(operational_analysis.get("risk_score", 0))],
        ["سطح ریسک", operational_analysis.get("risk_level") or "—"],
        ["درخواست باز", str(operational_analysis.get("open_requests", 0))],
        ["درخواست فوری", str(operational_analysis.get("urgent_requests", 0))],
        ["مسئله فوری/بحرانی", str(operational_analysis.get("critical_issues", 0))],
        ["اقدام معوق", str(operational_analysis.get("overdue_actions", 0))],
        ["شکاف خدماتی", str(operational_analysis.get("service_gap_count", 0))],
        ["تراکم مسئله در هکتار", str(operational_analysis.get("issue_density_per_ha", 0))],
    ]
    flowables.append(_make_pdf_table(["شاخص", "مقدار"], analysis_rows, col_widths=[75*mm, 95*mm]))
    flowables.append(Spacer(1, 6 * mm))

    # مسائل و نیازهای محله به ترتیب امتیاز
    flowables.append(Paragraph(shape_fa(f"مسائل و نیازهای محله (تعداد: {len(neighborhood_issues)})"), section_style))
    if neighborhood_issues:
        issue_rows = [[
            str(i), issue.get("title") or "—", issue.get("category") or "—",
            str(issue.get("priority_score") or 0), issue.get("priority_level") or "—",
            issue.get("status") or "—", issue.get("related_office") or "—"
        ] for i, issue in enumerate(neighborhood_issues, start=1)]
        flowables.append(_make_pdf_table(
            ["ردیف", "عنوان", "دسته", "امتیاز", "سطح", "وضعیت", "دستگاه"], issue_rows,
            col_widths=[12 * mm, 44 * mm, 27 * mm, 18 * mm, 20 * mm, 24 * mm, 30 * mm],
        ))
    else:
        flowables.append(Paragraph(shape_fa("مسئله‌ای برای این بلوک ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    flowables.append(Paragraph(shape_fa(f"اقدامات اجرایی (تعداد: {len(neighborhood_actions)})"), section_style))
    if neighborhood_actions:
        action_rows = [[
            str(i), action.get("title") or "—", action.get("responsible_office") or action.get("responsible_person") or "—",
            action.get("status") or "—", f"{action.get('progress_percent') or 0}٪", action.get("planned_end") or "—"
        ] for i, action in enumerate(neighborhood_actions, start=1)]
        flowables.append(_make_pdf_table(
            ["ردیف", "عنوان اقدام", "مسئول", "وضعیت", "پیشرفت", "پایان برنامه"], action_rows,
            col_widths=[12 * mm, 58 * mm, 36 * mm, 28 * mm, 18 * mm, 28 * mm],
        ))
    else:
        flowables.append(Paragraph(shape_fa("اقدام اجرایی ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    flowables.append(Paragraph(shape_fa(f"جلسات شورای محله (تعداد: {len(neighborhood_meetings)})"), section_style))
    if neighborhood_meetings:
        meeting_rows = [[
            str(i), item.get("title") or "—", item.get("meeting_date") or "—",
            item.get("place_name") or "—", item.get("status") or "—"
        ] for i, item in enumerate(neighborhood_meetings, start=1)]
        flowables.append(_make_pdf_table(
            ["ردیف", "عنوان", "تاریخ", "محل", "وضعیت"], meeting_rows,
            col_widths=[14 * mm, 58 * mm, 28 * mm, 48 * mm, 28 * mm],
        ))
    else:
        flowables.append(Paragraph(shape_fa("جلسه‌ای ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    flowables.append(Paragraph(shape_fa(f"مصوبات و پیگیری‌ها (تعداد: {len(neighborhood_resolutions)})"), section_style))
    if neighborhood_resolutions:
        resolution_rows = [[
            str(i), item.get("title") or "—", item.get("responsible_office") or item.get("responsible_person") or "—",
            item.get("due_date") or "—", item.get("status") or "—"
        ] for i, item in enumerate(neighborhood_resolutions, start=1)]
        flowables.append(_make_pdf_table(
            ["ردیف", "عنوان مصوبه", "مسئول", "مهلت", "وضعیت"], resolution_rows,
            col_widths=[14 * mm, 65 * mm, 42 * mm, 28 * mm, 31 * mm],
        ))
    else:
        flowables.append(Paragraph(shape_fa("مصوبه‌ای ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6 * mm))

    flowables.append(Paragraph(shape_fa("بودجه و هزینه بلوک"), section_style))
    budget_summary_rows = [
        ["اعتبار مصوب", f"{budget_summary.get('approved', 0):,.0f} ریال"],
        ["تخصیص‌یافته", f"{budget_summary.get('allocated', 0):,.0f} ریال"],
        ["هزینه‌شده", f"{budget_summary.get('spent', 0):,.0f} ریال"],
        ["مانده تخصیص", f"{budget_summary.get('remaining', 0):,.0f} ریال"],
        ["درصد مصرف", f"{budget_summary.get('utilization_percent', 0):.1f}٪"],
        ["تعداد اضافه‌هزینه", str(budget_summary.get('overrun_count', 0))],
    ]
    flowables.append(_make_pdf_table(["شاخص", "مقدار"], budget_summary_rows, col_widths=[70 * mm, 100 * mm]))
    if neighborhood_budgets:
        rows = [[str(i), b.get("title") or "—", b.get("funding_source") or "—",
                 f"{b.get('allocated_amount') or 0:,.0f}", f"{b.get('spent_amount') or 0:,.0f}",
                 b.get("status") or "—"] for i, b in enumerate(neighborhood_budgets, start=1)]
        flowables.append(_make_pdf_table(
            ["ردیف", "عنوان", "منبع", "تخصیص", "هزینه", "وضعیت"], rows,
            col_widths=[12 * mm, 50 * mm, 35 * mm, 27 * mm, 27 * mm, 29 * mm],
        ))
    flowables.append(Spacer(1, 6 * mm))

    flowables.append(Paragraph(shape_fa("ارزیابی عملکرد و کنترل کیفیت"), section_style))
    performance_rows = [
        ["امتیاز کل", f"{performance.get('total_score', 0):.1f} از ۱۰۰"],
        ["سطح عملکرد", performance.get("level") or "—"],
        ["تکمیل اطلاعات", f"{performance.get('completeness', 0):.1f}٪"],
        ["حل مسائل", f"{performance.get('issue_resolution', 0):.1f}٪"],
        ["تکمیل اقدامات", f"{performance.get('action_completion', 0):.1f}٪"],
        ["تحقق مصوبات", f"{performance.get('resolution_completion', 0):.1f}٪"],
        ["رعایت زمان‌بندی", f"{performance.get('timeliness', 0):.1f}٪"],
        ["کنترل مالی", f"{performance.get('financial_control', 0):.1f}٪"],
        ["پاسخ‌گویی و مشارکت", f"{performance.get('participation_response', 0):.1f}٪"],
    ]
    flowables.append(_make_pdf_table(["شاخص", "مقدار"], performance_rows, col_widths=[70 * mm, 100 * mm]))
    if management_alerts:
        alert_rows = [[a.get("severity"), a.get("category"), a.get("title"), a.get("due_date") or "—"]
                      for a in management_alerts[:20]]
        flowables.append(Paragraph(shape_fa(f"هشدارهای باز (تعداد: {len(management_alerts)})"), section_style))
        flowables.append(_make_pdf_table(["سطح", "دسته", "عنوان", "مهلت"], alert_rows,
                                         col_widths=[25 * mm, 30 * mm, 90 * mm, 30 * mm]))
    if quality_issues:
        quality_rows = [[q.get("severity"), q.get("category"), q.get("message")] for q in quality_issues[:20]]
        flowables.append(Paragraph(shape_fa(f"مغایرت‌های پرونده (تعداد: {len(quality_issues)})"), section_style))
        flowables.append(_make_pdf_table(["شدت", "دسته", "شرح"], quality_rows,
                                         col_widths=[25 * mm, 35 * mm, 115 * mm]))

    doc.build(flowables)
    return output_path

def generate_block_full_report_excel(db, zone_id, output_path, map_image_path=None):
    """نسخه Excel گزارش کامل بلوک؛ تصویر در ابتدای شیت اول و سپس مشخصات کامل."""
    zone = db.get_zone(zone_id)
    if not zone:
        raise ValueError("منطقه یافت نشد")

    streets = db.get_streets(zone_id=zone_id)
    places = db.get_places(zone_id=zone_id)
    mosques = db.get_mosques(zone_id=zone_id)
    members = db.get_council_members(zone_id=zone_id)
    requests = db.get_priority_requests(zone_id=zone_id)
    meeting_place = db.get_zone_meeting_place(zone_id)
    profile = db.get_zone_profile(zone_id)
    neighborhood_issues = db.get_neighborhood_issues(zone_id)
    neighborhood_actions = db.get_neighborhood_actions(zone_id)
    neighborhood_meetings = db.get_neighborhood_meetings(zone_id)
    neighborhood_resolutions = db.get_neighborhood_resolutions(zone_id=zone_id)
    neighborhood_budgets = db.get_neighborhood_budgets(zone_id)
    budget_summary = db.get_budget_summary(zone_id)
    performance = db.get_zone_performance(zone_id)
    management_alerts = db.get_management_alerts(zone_id)
    quality_issues = db.get_quality_issues(zone_id)
    field_visits = db.get_field_visits(zone_id)
    citizen_requests = db.get_citizen_requests(zone_id)
    operational_analysis = db.get_zone_operational_analysis(zone_id)
    correspondence_letters = db.get_correspondence_letters(zone_id=zone_id)
    approval_requests = db.get_approval_requests(zone_id=zone_id, limit=5000)
    generated_documents = db.get_generated_documents(zone_id=zone_id, limit=5000)
    committees = get_zone_committee_report_data(db, zone_id)

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "گزارش کامل بلوک"
    ws.sheet_view.rightToLeft = True
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_setup.orientation = "portrait"
    ws.page_margins.left = 0.35
    ws.page_margins.right = 0.35
    ws.page_margins.top = 0.45
    ws.page_margins.bottom = 0.45

    ws.merge_cells("A1:D1")
    ws["A1"] = f"گزارش کامل بلوک: {zone['name']}"
    ws["A1"].font = TITLE_FONT
    ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 24

    image_source, snapshot_meta = _snapshot_image_source(db, zone_id, map_image_path)
    if image_source is not None:
        try:
            from openpyxl.drawing.image import Image as XLImage
            xl_img = XLImage(image_source)
            # تقریباً نیم صفحه A4 در حالت عمودی.
            xl_img.width = 620
            xl_img.height = 465
            ws.add_image(xl_img, "A3")
        except Exception:
            ws["A3"] = "خطا در بارگذاری تصویر گرافیکی بلوک."
    else:
        ws["A3"] = "نمای گرافیکی بلوک در دسترس نیست."

    summary_start = 29
    ws.cell(row=summary_start, column=1, value="مشخصات پایه بلوک").font = HEADER_FONT
    summary_rows = [
        ["نام بلوک", zone.get("name") or "—"],
        ["وضعیت تکمیل", zone.get("status") or "—"],
        ["مساحت", f"{(zone.get('area_m2', 0) or 0)/10000:.2f} هکتار"],
        ["محیط", f"{zone.get('perimeter_m', 0) or 0:.0f} متر"],
        ["تعداد نقاط مرزی", len(zone.get("boundary_points", []))],
        ["تعداد خیابان و کوچه", len(streets)],
        ["تعداد سایر اماکن", len(places)],
        ["تعداد مساجد", len(mosques)],
        ["تعداد کمیته‌های تخصصی", len(committees)],
        ["اعضای فعال کمیته‌ها", sum(len(c.get("active_members") or []) for c in committees)],
        ["مصوبات باز کمیته‌ها", sum(1 for c in committees for r in (c.get("resolutions") or []) if r.get("status") not in ("انجام‌شده", "لغوشده"))],
        ["بازدیدهای میدانی", len(field_visits)],
        ["درخواست‌های مردمی", len(citizen_requests)],
        ["مکاتبات اداری مرتبط", len(correspondence_letters)],
        ["گردش‌های تأیید", len(approval_requests)],
        ["اسناد Word تولیدشده", len(generated_documents)],
        ["سطح ریسک عملیاتی", operational_analysis.get("risk_level") or "—"],
        ["خانوار نهایی تأییدشده", profile.get("approved_households", 0)],
        ["جمعیت تخمینی", profile.get("estimated_population", 0)],
        ["مسائل باز", sum(1 for i in neighborhood_issues if i.get("status") not in ("مختومه", "انجام‌شده"))],
        ["اقدامات اجرایی", len(neighborhood_actions)],
        ["مصوبات باز", sum(1 for r in neighborhood_resolutions if r.get("status") not in ("انجام‌شده", "لغوشده"))],
        ["تاریخ ایجاد", zone.get("created_at") or "—"],
        ["آخرین بروزرسانی", zone.get("updated_at") or "—"],
    ]
    current_row = summary_start + 1
    for label, value in summary_rows:
        ws.cell(row=current_row, column=1, value=label)
        ws.cell(row=current_row, column=2, value=value)
        current_row += 1

    current_row += 1
    ws.cell(row=current_row, column=1, value="محل برگزاری جلسات").font = HEADER_FONT
    current_row += 1
    ws.cell(row=current_row, column=1, value="نام مکان")
    ws.cell(row=current_row, column=2, value=meeting_place["place_name"] if meeting_place else "—")
    current_row += 1
    ws.cell(row=current_row, column=1, value="آدرس دقیق")
    ws.cell(row=current_row, column=2, value=meeting_place["exact_address"] if meeting_place else "—")

    ws.column_dimensions["A"].width = 25
    ws.column_dimensions["B"].width = 55
    ws.column_dimensions["C"].width = 20
    ws.column_dimensions["D"].width = 20
    for row in ws.iter_rows(min_row=summary_start, max_row=current_row, min_col=1, max_col=2):
        for cell in row:
            cell.alignment = Alignment(horizontal="right", vertical="center", wrap_text=True)

    # شیت‌های جزئیات به ترتیب منطقی گزارش.
    ws_streets = wb.create_sheet("خیابان‌ها و کوچه‌ها")
    _write_excel_sheet(
        ws_streets, f"خیابان‌ها و کوچه‌های بلوک {zone['name']}",
        ["ردیف", "نام معبر", "نوع"],
        [[i, s.get("name") or "معبر بدون نام", s.get("highway_type") or ""]
         for i, s in enumerate(streets, start=1)],
    )

    ws_places = wb.create_sheet("اماکن")
    _write_excel_sheet(
        ws_places, f"سایر اماکن بلوک {zone['name']}",
        ["ردیف", "نام مکان", "دسته", "نوع", "عرض", "طول"],
        [[i, p.get("name") or "", p.get("category") or "", p.get("subtype") or "", p.get("lat"), p.get("lon")]
         for i, p in enumerate(places, start=1)],
    )

    ws_mosques = wb.create_sheet("مساجد")
    _write_excel_sheet(
        ws_mosques, f"مساجد بلوک {zone['name']}",
        ["ردیف", "نام مسجد", "عرض جغرافیایی", "طول جغرافیایی"],
        [[i, m["name"], m["lat"], m["lon"]] for i, m in enumerate(mosques, start=1)],
    )

    ws_members = wb.create_sheet("معتمدین")
    _write_excel_sheet(
        ws_members, f"معتمدین بلوک {zone['name']}",
        ["نام و نام‌خانوادگی", "دسته", "سمت", "شماره تماس"],
        [[f"{m['first_name']} {m['last_name']}", m["member_group"] or "", m["position"] or "", m["mobile"] or ""]
         for m in members],
    )

    ws_committees = wb.create_sheet("کمیته‌های شش‌گانه")
    _write_excel_sheet(
        ws_committees, f"کمیته‌های تخصصی بلوک {zone['name']}",
        ["ردیف", "نام کمیته", "رئیس", "دبیر", "اعضای فعال", "جلسات", "مصوبات باز", "وضعیت"],
        committee_summary_rows(committees),
    )

    ws_committee_members = wb.create_sheet("اعضای کمیته‌ها")
    _write_excel_sheet(
        ws_committee_members, f"اعضا و سمت‌های کمیته‌های بلوک {zone['name']}",
        ["کمیته", "نام عضو", "کد ملی", "سمت", "نوع عضویت", "اداره/دستگاه", "موبایل", "وضعیت"],
        committee_member_rows(committees),
    )

    ws_committee_meetings = wb.create_sheet("جلسات کمیته‌ها")
    _write_excel_sheet(
        ws_committee_meetings, f"جلسات کمیته‌های بلوک {zone['name']}",
        ["کمیته", "عنوان جلسه", "تاریخ", "ساعت", "محل", "وضعیت"],
        committee_meeting_rows(committees),
    )

    ws_committee_resolutions = wb.create_sheet("مصوبات کمیته‌ها")
    _write_excel_sheet(
        ws_committee_resolutions, f"مصوبات کمیته‌های بلوک {zone['name']}",
        ["کمیته", "عنوان مصوبه", "مسئول پیگیری", "دستگاه مسئول", "مهلت", "وضعیت", "مسئله مرتبط", "اقدام مرتبط"],
        committee_resolution_rows(committees),
    )

    ws_committee_links = wb.create_sheet("ارتباطات کمیته‌ها")
    _write_excel_sheet(
        ws_committee_links, f"مسائل و اقدامات مرتبط با کمیته‌های بلوک {zone['name']}",
        ["کمیته", "نوع", "عنوان", "دستگاه", "وضعیت"],
        committee_link_rows(committees),
    )

    ws_profile = wb.create_sheet("جمعیت و خانوار")
    _write_excel_sheet(
        ws_profile, f"جمعیت و خانوار بلوک {zone['name']}",
        ["شاخص", "مقدار"],
        [
            ["ساختمان مسکونی", profile.get("residential_buildings", 0)],
            ["واحد مسکونی", profile.get("residential_units", 0)],
            ["واحد دارای سکنه", profile.get("occupied_units", 0)],
            ["واحد خالی", profile.get("vacant_units", 0)],
            ["خانوار محاسبه‌شده", profile.get("estimated_households", 0)],
            ["خانوار ثبت میدانی", profile.get("field_households", 0)],
            ["خانوار نهایی تأییدشده", profile.get("approved_households", 0)],
            ["جمعیت تخمینی", profile.get("estimated_population", 0)],
            ["روش برآورد", profile.get("estimation_method") or ""],
            ["سطح اطمینان", profile.get("confidence_level") or ""],
        ],
    )

    ws_issues = wb.create_sheet("مسائل و نیازها")
    _write_excel_sheet(
        ws_issues, f"مسائل بلوک {zone['name']}",
        ["شناسه", "عنوان", "دسته", "شرح", "دستگاه", "امتیاز", "سطح", "وضعیت", "خانوار تحت تأثیر", "مهلت"],
        [[i.get("id"), i.get("title"), i.get("category"), i.get("description"), i.get("related_office"),
          i.get("priority_score"), i.get("priority_level"), i.get("status"), i.get("affected_households"), i.get("due_date")]
         for i in neighborhood_issues],
    )

    ws_actions = wb.create_sheet("اقدامات اجرایی")
    _write_excel_sheet(
        ws_actions, f"اقدامات بلوک {zone['name']}",
        ["شناسه", "عنوان", "مسئله مرتبط", "مسئول", "دستگاه", "شروع", "پایان", "پیشرفت", "وضعیت", "وضعیت تأیید", "هزینه برآوردی", "هزینه واقعی"],
        [[a.get("id"), a.get("title"), a.get("issue_title"), a.get("responsible_person"), a.get("responsible_office"),
          a.get("planned_start"), a.get("planned_end"), a.get("progress_percent"), a.get("status"), a.get("approval_status"),
          a.get("estimated_cost"), a.get("actual_cost")] for a in neighborhood_actions],
    )

    ws_meetings = wb.create_sheet("جلسات")
    _write_excel_sheet(
        ws_meetings, f"جلسات بلوک {zone['name']}",
        ["شناسه", "عنوان", "تاریخ", "ساعت", "محل", "دستور جلسه", "حاضرین", "صورت‌جلسه", "وضعیت"],
        [[m.get("id"), m.get("title"), m.get("meeting_date"), m.get("start_time"), m.get("place_name"),
          m.get("agenda"), m.get("attendees"), m.get("minutes_text"), m.get("status")] for m in neighborhood_meetings],
    )

    ws_resolutions = wb.create_sheet("مصوبات")
    _write_excel_sheet(
        ws_resolutions, f"مصوبات بلوک {zone['name']}",
        ["شناسه", "جلسه", "عنوان", "شرح", "مسئول", "دستگاه", "مهلت", "وضعیت", "وضعیت تأیید", "مسئله مرتبط", "اقدام مرتبط"],
        [[r.get("id"), r.get("meeting_title"), r.get("title"), r.get("description"), r.get("responsible_person"),
          r.get("responsible_office"), r.get("due_date"), r.get("status"), r.get("approval_status"), r.get("linked_issue_id"), r.get("linked_action_id")]
         for r in neighborhood_resolutions],
    )

    ws_letters = wb.create_sheet("مکاتبات اداری")
    _write_excel_sheet(
        ws_letters, f"مکاتبات مرتبط با بلوک {zone['name']}",
        ["شناسه", "شماره", "نوع", "موضوع", "فرستنده", "گیرنده", "تاریخ نامه", "مهلت", "وضعیت", "وضعیت تأیید", "اولویت", "طبقه‌بندی", "پیوست", "ارجاع باز"],
        [[l.get("id"), l.get("letter_number"), l.get("direction"), l.get("subject"), l.get("sender"),
          l.get("recipient"), l.get("letter_date"), l.get("due_date"), l.get("status"), l.get("approval_status"), l.get("priority"),
          l.get("confidentiality"), l.get("attachment_count"), l.get("open_assignment_count")]
         for l in correspondence_letters],
    )

    ws_approvals = wb.create_sheet("گردش تأیید")
    _write_excel_sheet(
        ws_approvals, f"گردش‌های تأیید بلوک {zone['name']}",
        ["شناسه", "عنوان", "نوع رکورد", "شناسه رکورد", "مرحله", "تعداد مراحل", "مهلت", "وضعیت", "درخواست‌کننده", "توضیحات"],
        [[a.get("id"), a.get("title"), a.get("entity_type"), a.get("entity_id"),
          a.get("current_step"), a.get("total_steps"), a.get("due_date"), a.get("status"),
          a.get("requested_by_name"), a.get("notes")] for a in approval_requests],
    )

    ws_generated = wb.create_sheet("اسناد تولیدشده")
    _write_excel_sheet(
        ws_generated, f"اسناد Word تولیدشده برای بلوک {zone['name']}",
        ["شناسه", "عنوان", "قالب", "نوع ارتباط", "شناسه ارتباط", "تولیدکننده", "تاریخ", "مسیر فایل"],
        [[d.get("id"), d.get("title"), d.get("template_name"), d.get("related_entity_type"),
          d.get("related_entity_id"), d.get("created_by_name"), d.get("created_at"), d.get("file_path")]
         for d in generated_documents],
    )

    ws_visits = wb.create_sheet("بازدیدهای میدانی")
    _write_excel_sheet(
        ws_visits, f"بازدیدهای میدانی بلوک {zone['name']}",
        ["شناسه", "تاریخ", "ساعت", "کارشناس", "نوع", "موقعیت", "عرض", "طول", "ساختمان", "خانوار", "مشاهدات", "اقدام فوری", "پیگیری", "وضعیت"],
        [[v.get("id"), v.get("visit_date"), v.get("start_time"), v.get("officer_name"), v.get("visit_type"),
          v.get("location_text"), v.get("lat"), v.get("lon"), v.get("buildings_count"), v.get("households_count"),
          v.get("observation"), v.get("immediate_action"), "بله" if v.get("followup_required") else "خیر", v.get("status")]
         for v in field_visits],
    )

    ws_citizen = wb.create_sheet("درخواست‌های مردمی")
    _write_excel_sheet(
        ws_citizen, f"درخواست‌های مردمی بلوک {zone['name']}",
        ["کد رهگیری", "نام شهروند", "ناشناس", "موبایل", "دسته", "عنوان", "شرح", "موقعیت", "فوریت", "دستگاه", "وضعیت", "مسئله مرتبط", "منبع", "تاریخ دریافت"],
        [[r.get("tracking_code"), r.get("citizen_name"), "بله" if r.get("is_anonymous") else "خیر", r.get("mobile"),
          r.get("category"), r.get("title"), r.get("description"), r.get("location_text"), r.get("urgency"),
          r.get("assigned_office"), r.get("status"), r.get("linked_issue_id"), r.get("source"), r.get("received_at")]
         for r in citizen_requests],
    )

    ws_operational = wb.create_sheet("تحلیل عملیاتی")
    _write_excel_sheet(
        ws_operational, f"تحلیل عملیاتی بلوک {zone['name']}",
        ["شاخص", "مقدار"],
        [["مساحت هکتار", operational_analysis.get("area_ha")],
         ["امتیاز ریسک", operational_analysis.get("risk_score")],
         ["سطح ریسک", operational_analysis.get("risk_level")],
         ["بازدید میدانی", operational_analysis.get("field_visits")],
         ["درخواست مردمی", operational_analysis.get("citizen_requests")],
         ["درخواست باز", operational_analysis.get("open_requests")],
         ["درخواست فوری", operational_analysis.get("urgent_requests")],
         ["مسئله بحرانی", operational_analysis.get("critical_issues")],
         ["اقدام معوق", operational_analysis.get("overdue_actions")],
         ["شکاف خدماتی", operational_analysis.get("service_gap_count")],
         ["تراکم مسئله در هکتار", operational_analysis.get("issue_density_per_ha")],
         ["تراکم درخواست در هکتار", operational_analysis.get("request_density_per_ha")]],
    )

    ws_budget = wb.create_sheet("بودجه و هزینه")
    _write_excel_sheet(
        ws_budget, f"بودجه بلوک {zone['name']}",
        ["شناسه", "عنوان", "اقدام مرتبط", "سال", "منبع", "مصوب", "تخصیص", "هزینه", "مانده", "وضعیت", "وضعیت تأیید", "سند"],
        [[b.get("id"), b.get("title"), b.get("action_title"), b.get("fiscal_year"), b.get("funding_source"),
          b.get("approved_amount"), b.get("allocated_amount"), b.get("spent_amount"),
          (b.get("allocated_amount") or 0) - (b.get("spent_amount") or 0), b.get("status"), b.get("approval_status"), b.get("document_reference")]
         for b in neighborhood_budgets],
    )

    ws_performance = wb.create_sheet("ارزیابی عملکرد")
    _write_excel_sheet(
        ws_performance, f"ارزیابی عملکرد بلوک {zone['name']}",
        ["شاخص", "مقدار"],
        [["امتیاز کل", performance.get("total_score")], ["سطح", performance.get("level")],
         ["تکمیل اطلاعات", performance.get("completeness")], ["حل مسائل", performance.get("issue_resolution")],
         ["تکمیل اقدامات", performance.get("action_completion")], ["تحقق مصوبات", performance.get("resolution_completion")],
         ["رعایت زمان‌بندی", performance.get("timeliness")], ["کنترل مالی", performance.get("financial_control")],
         ["اعتبار مصوب", budget_summary.get("approved")], ["تخصیص", budget_summary.get("allocated")],
         ["هزینه", budget_summary.get("spent")], ["مانده", budget_summary.get("remaining")]],
    )

    ws_alerts = wb.create_sheet("هشدارها")
    _write_excel_sheet(
        ws_alerts, f"هشدارهای باز بلوک {zone['name']}",
        ["سطح", "دسته", "عنوان", "توضیح", "مهلت", "نوع موجودیت", "شناسه"],
        [[a.get("severity"), a.get("category"), a.get("title"), a.get("detail"), a.get("due_date"),
          a.get("entity_type"), a.get("entity_id")] for a in management_alerts],
    )

    ws_quality = wb.create_sheet("کنترل کیفیت")
    _write_excel_sheet(
        ws_quality, f"مغایرت‌های پرونده بلوک {zone['name']}",
        ["شدت", "دسته", "شرح"],
        [[q.get("severity"), q.get("category"), q.get("message")] for q in quality_issues],
    )

    wb.save(output_path)
    return output_path


def _trusted_members_for_public_report(db, zone_id):
    """فقط اعضایی که دسته آن‌ها «معتمد» است؛ با حفظ ترتیب ثبت."""
    members = db.get_council_members(zone_id=zone_id)
    return [member for member in members if "معتمد" in str(member.get("member_group") or "").strip()]


def _fit_canvas_text(text, font_name, font_size, max_width):
    """کوتاه‌سازی امن متن برای جلوگیری از خروج از سلول جدول ثابت A4."""
    value = shape_fa(str(text if text not in (None, "") else "—"))
    if pdfmetrics.stringWidth(value, font_name, font_size) <= max_width:
        return value
    raw = str(text if text not in (None, "") else "—")
    while raw:
        raw = raw[:-1]
        candidate = shape_fa(raw + "...")
        if pdfmetrics.stringWidth(candidate, font_name, font_size) <= max_width:
            return candidate
    return "..."


def generate_block_public_report_pdf(db, zone_id, output_path, map_image_path=None):
    """گزارش عمومی تک‌صفحه‌ای A4: نقشه در یک‌چهارم بالا و جدول معتمدین در ادامه."""
    zone = db.get_zone(zone_id)
    if not zone:
        raise ValueError("منطقه یافت نشد")

    _ensure_fonts_registered()
    page_width, page_height = A4
    pdf = canvas.Canvas(output_path, pagesize=A4)
    pdf.setTitle(f"Public block report - {zone.get('name') or zone_id}")
    margin_x = 12 * mm
    content_width = page_width - (2 * margin_x)
    top_y = page_height - 11 * mm

    pdf.setFillColor(colors.HexColor("#13294b"))
    pdf.setFont(FONT_NAME_BOLD, 12)
    pdf.drawCentredString(page_width / 2, top_y, shape_fa("فرمانداری شهرستان جوانرود"))
    pdf.setFont(FONT_NAME_BOLD, 15)
    pdf.drawCentredString(page_width / 2, top_y - 8 * mm, shape_fa(f"گزارش عمومی بلوک: {zone.get('name') or '—'}"))
    pdf.setFont(FONT_NAME, 8)
    pdf.setFillColor(colors.HexColor("#5b6472"))
    pdf.drawRightString(page_width - margin_x, top_y - 14 * mm, shape_fa(f"تاریخ تهیه: {now_jalali()}"))
    pdf.setStrokeColor(colors.HexColor("#c9a227"))
    pdf.setLineWidth(1.2)
    pdf.line(margin_x, top_y - 17 * mm, page_width - margin_x, top_y - 17 * mm)

    map_top = top_y - 21 * mm
    map_height = 62 * mm
    map_bottom = map_top - map_height
    pdf.setFillColor(colors.white)
    pdf.setStrokeColor(colors.HexColor("#cbd2dc"))
    pdf.roundRect(margin_x, map_bottom, content_width, map_height, 3 * mm, stroke=1, fill=1)

    image_source, _snapshot_meta = _snapshot_image_source(db, zone_id, map_image_path)
    if image_source is not None:
        try:
            reader = ImageReader(image_source)
            src_w, src_h = reader.getSize()
            inner_pad = 2.5 * mm
            box_w = content_width - 2 * inner_pad
            box_h = map_height - 2 * inner_pad
            scale = min(box_w / float(src_w), box_h / float(src_h))
            draw_w, draw_h = float(src_w) * scale, float(src_h) * scale
            draw_x = margin_x + (content_width - draw_w) / 2
            draw_y = map_bottom + (map_height - draw_h) / 2
            pdf.drawImage(reader, draw_x, draw_y, draw_w, draw_h, preserveAspectRatio=True, mask='auto')
        except Exception:
            pdf.setFillColor(colors.HexColor("#5b6472"))
            pdf.setFont(FONT_NAME, 9)
            pdf.drawCentredString(page_width / 2, map_bottom + map_height / 2, shape_fa("نمای نقشه بلوک قابل بارگذاری نیست."))
    else:
        pdf.setFillColor(colors.HexColor("#5b6472"))
        pdf.setFont(FONT_NAME, 9)
        pdf.drawCentredString(page_width / 2, map_bottom + map_height / 2, shape_fa("نمای نقشه بلوک هنوز تولید نشده است."))

    members = _trusted_members_for_public_report(db, zone_id)
    table_title_y = map_bottom - 8 * mm
    pdf.setFillColor(colors.HexColor("#13294b"))
    pdf.setFont(FONT_NAME_BOLD, 11)
    pdf.drawRightString(page_width - margin_x, table_title_y, shape_fa(f"جدول اعضای معتمد - تعداد: {len(members)}"))

    table_top = table_title_y - 4 * mm
    table_bottom = 13 * mm
    available_height = table_top - table_bottom
    row_count = max(1, len(members)) + 1
    row_height = min(8 * mm, available_height / row_count)
    font_size = max(6.2, min(9.0, (row_height / mm) * 1.05))
    padding = 2.2 * mm

    widths = [42 * mm, 32 * mm, 32 * mm, 63 * mm, 15 * mm]
    x_positions = [margin_x]
    for width in widths:
        x_positions.append(x_positions[-1] + width)

    table_rows = [["سمت", "شماره تماس", "کد ملی", "نام و نام خانوادگی", "ردیف"]]
    for index, member in enumerate(members, start=1):
        full_name = f"{member.get('first_name') or ''} {member.get('last_name') or ''}".strip() or "—"
        table_rows.append([member.get("position") or "—", member.get("mobile") or "—", member.get("national_code") or "—", full_name, str(index)])
    if not members:
        table_rows.append(["—", "—", "—", "عضو معتمدی برای این بلوک ثبت نشده است.", "—"])
        row_height = min(8 * mm, available_height / 2)
        font_size = 8.5

    y = table_top
    for row_index, row in enumerate(table_rows):
        y_next = y - row_height
        pdf.setFillColor(colors.HexColor("#13294b") if row_index == 0 else (colors.HexColor("#f0f2f6") if row_index % 2 == 0 else colors.white))
        pdf.rect(margin_x, y_next, sum(widths), row_height, stroke=0, fill=1)
        pdf.setStrokeColor(colors.HexColor("#d7dbe3"))
        pdf.setLineWidth(0.45)
        for x in x_positions:
            pdf.line(x, y_next, x, y)
        pdf.line(margin_x, y, margin_x + sum(widths), y)
        pdf.line(margin_x, y_next, margin_x + sum(widths), y_next)
        font_name = FONT_NAME_BOLD if row_index == 0 else FONT_NAME
        pdf.setFillColor(colors.white if row_index == 0 else colors.HexColor("#1c2530"))
        pdf.setFont(font_name, font_size)
        baseline = y_next + (row_height - font_size) / 2 + 1.2
        for col_index, cell in enumerate(row):
            left, right = x_positions[col_index], x_positions[col_index + 1]
            fitted = _fit_canvas_text(cell, font_name, font_size, (right - left) - 2 * padding)
            pdf.drawCentredString((left + right) / 2, baseline, fitted)
        y = y_next

    pdf.setFillColor(colors.HexColor("#7a8491"))
    pdf.setFont(FONT_NAME, 7)
    pdf.drawCentredString(page_width / 2, 7 * mm, shape_fa("سامانه مدیریت محلات جوانرود"))
    pdf.showPage()
    pdf.save()
    return output_path

def generate_zone_full_report_pdf(db, zone_id, output_path):
    """گزارش کامل منطقه؛ از همان قالب جامع بلوک استفاده می‌کند."""
    return generate_block_full_report_pdf(db, zone_id, output_path)

def generate_members_report_pdf(db, output_path, zone_id=None):
    """گزارش اعضای شورای محلات (همه مناطق یا یک منطقه خاص)."""
    doc = SimpleDocTemplate(output_path, pagesize=A4,
                             topMargin=15 * mm, bottomMargin=15 * mm,
                             leftMargin=15 * mm, rightMargin=15 * mm)
    title = "گزارش اعضای شورای محلات"
    zone_name = None
    if zone_id is not None:
        zone = db.get_zone(zone_id)
        if zone:
            zone_name = zone["name"]
            title += f" — منطقه: {zone_name}"

    flowables = _pdf_header_flowables(title)
    members = db.get_council_members(zone_id=zone_id)

    zones_by_id = {z["id"]: z["name"] for z in db.get_zones()}
    headers = ["نام", "نام خانوادگی", "کد ملی", "تحصیلات", "موبایل", "دسته", "سمت"]
    if zone_id is None:
        headers.append("منطقه")

    rows = []
    for m in members:
        row = [m["first_name"], m["last_name"], m["national_code"] or "—",
               m["education"] or "—", m["mobile"] or "—", m["member_group"] or "—", m["position"] or "—"]
        if zone_id is None:
            row.append(zones_by_id.get(m["zone_id"], "—"))
        rows.append(row)

    if rows:
        col_widths = [22 * mm, 22 * mm, 22 * mm, 22 * mm, 25 * mm, 20 * mm, 25 * mm]
        if zone_id is None:
            col_widths.append(22 * mm)
        flowables.append(_make_pdf_table(headers, rows, col_widths=col_widths))
    else:
        _, _, _, normal_style = _pdf_styles()
        flowables.append(Paragraph(shape_fa("هیچ عضوی ثبت نشده است."), normal_style))

    doc.build(flowables)
    return output_path


def generate_requests_report_pdf(db, output_path, zone_id=None):
    """گزارش درخواست‌ها و مشکلات اولویت‌بندی‌شده."""
    doc = SimpleDocTemplate(output_path, pagesize=A4,
                             topMargin=15 * mm, bottomMargin=15 * mm,
                             leftMargin=15 * mm, rightMargin=15 * mm)
    title = "گزارش درخواست‌ها و مشکلات اولویت‌بندی‌شده"
    if zone_id is not None:
        zone = db.get_zone(zone_id)
        if zone:
            title += f" — منطقه: {zone['name']}"

    flowables = _pdf_header_flowables(title)
    requests = db.get_priority_requests(zone_id=zone_id)
    zones_by_id = {z["id"]: z["name"] for z in db.get_zones()}

    headers = ["ردیف", "منطقه", "شرح درخواست/مشکل", "اداره مرتبط", "تعداد اقدام"]
    rows = []
    for i, r in enumerate(requests, start=1):
        rows.append([
            str(i),
            zones_by_id.get(r["zone_id"], "—"),
            r["description"],
            r["related_office"] or "—",
            str(r["action_count"])
        ])

    if rows:
        flowables.append(_make_pdf_table(headers, rows,
                                          col_widths=[12 * mm, 30 * mm, 75 * mm, 35 * mm, 18 * mm]))
    else:
        _, _, _, normal_style = _pdf_styles()
        flowables.append(Paragraph(shape_fa("درخواستی ثبت نشده است."), normal_style))

    doc.build(flowables)
    return output_path


def generate_actions_report_pdf(db, output_path, zone_id=None):
    """گزارش کامل اقدامات انجام‌شده (تاریخچه پیگیری هر درخواست)."""
    doc = SimpleDocTemplate(output_path, pagesize=A4,
                             topMargin=15 * mm, bottomMargin=15 * mm,
                             leftMargin=15 * mm, rightMargin=15 * mm)
    title = "گزارش اقدامات انجام‌شده"
    if zone_id is not None:
        zone = db.get_zone(zone_id)
        if zone:
            title += f" — منطقه: {zone['name']}"

    flowables = _pdf_header_flowables(title)
    _, _, section_style, normal_style = _pdf_styles()
    requests = db.get_priority_requests(zone_id=zone_id)
    zones_by_id = {z["id"]: z["name"] for z in db.get_zones()}

    any_action_found = False
    for r in requests:
        actions = db.get_request_actions(r["id"])
        if not actions:
            continue
        any_action_found = True
        header_text = f"{zones_by_id.get(r['zone_id'], '—')} — {r['description']}"
        flowables.append(Paragraph(shape_fa(header_text), section_style))
        rows = [[a["created_at"], a["action_description"]] for a in actions]
        flowables.append(_make_pdf_table(["تاریخ", "شرح اقدام"], rows, col_widths=[35 * mm, 135 * mm]))
        flowables.append(Spacer(1, 4 * mm))

    if not any_action_found:
        flowables.append(Paragraph(shape_fa("هیچ اقدامی ثبت نشده است."), normal_style))

    doc.build(flowables)
    return output_path


# ==================== Excel ====================

HEADER_FILL = PatternFill(start_color="13294B", end_color="13294B", fill_type="solid")
HEADER_FONT = Font(color="FFFFFF", bold=True, size=11)
TITLE_FONT = Font(color="0B1F3A", bold=True, size=14)


def _write_excel_sheet(ws, title, headers, rows):
    ws.sheet_view.rightToLeft = True
    ws["A1"] = title
    ws["A1"].font = TITLE_FONT
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 1))

    header_row_idx = 3
    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=header_row_idx, column=col_idx, value=header)
        cell.font = HEADER_FONT
        cell.fill = HEADER_FILL
        cell.alignment = Alignment(horizontal="center", vertical="center")

    for row_offset, row_data in enumerate(rows, start=1):
        for col_idx, value in enumerate(row_data, start=1):
            cell = ws.cell(row=header_row_idx + row_offset, column=col_idx, value=value)
            cell.alignment = Alignment(horizontal="right", vertical="center", wrap_text=True)

    for col_idx in range(1, len(headers) + 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = 22


def generate_overview_report_excel(db, output_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "خلاصه وضعیت"

    zones = db.get_zones()
    all_streets = db.get_streets()
    all_places = db.get_places()
    all_members = db.get_council_members()
    all_requests = db.get_priority_requests()
    all_committees = [c for z in zones for c in get_zone_committee_report_data(db, z["id"])]
    total_actions = sum(r["action_count"] for r in all_requests)

    rows = [
        ["تعداد مناطق/بلوک‌ها", len(zones)],
        ["تعداد خیابان‌ها و کوچه‌ها", len(all_streets)],
        ["تعداد اماکن", len(all_places)],
        ["تعداد اعضای شورا", len(all_members)],
        ["تعداد کمیته‌های تخصصی", len(all_committees)],
        ["تعداد اعضای فعال کمیته‌ها", sum(len(c.get("active_members") or []) for c in all_committees)],
        ["مصوبات باز کمیته‌ها", sum(1 for c in all_committees for r in (c.get("resolutions") or []) if r.get("status") not in ("انجام‌شده", "لغوشده"))],
        ["تعداد درخواست‌ها", len(all_requests)],
        ["مجموع اقدامات پیگیری", total_actions],
    ]
    _write_excel_sheet(ws, "گزارش کلی وضعیت سامانه", ["شاخص", "مقدار"], rows)

    ws2 = wb.create_sheet("جزئیات مناطق")
    zone_rows = []
    for z in zones:
        zone_rows.append([
            z["name"],
            len(db.get_streets(zone_id=z["id"])),
            len(db.get_places(zone_id=z["id"])),
            len(db.get_council_members(zone_id=z["id"])),
            len(db.get_priority_requests(zone_id=z["id"])),
        ])
    _write_excel_sheet(ws2, "جزئیات هر منطقه",
                        ["نام منطقه", "تعداد خیابان", "تعداد مکان", "تعداد اعضا", "تعداد درخواست"], zone_rows)

    wb.save(output_path)
    return output_path


def generate_zone_full_report_excel(db, zone_id, output_path):
    """خروجی Excel کامل منطقه؛ از قالب جامع بلوک استفاده می‌کند."""
    return generate_block_full_report_excel(db, zone_id, output_path)

def generate_members_report_excel(db, output_path, zone_id=None):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "اعضای شورا"

    members = db.get_council_members(zone_id=zone_id)
    zones_by_id = {z["id"]: z["name"] for z in db.get_zones()}

    headers = ["نام", "نام خانوادگی", "کد ملی", "تحصیلات", "موبایل", "دسته", "سمت"]
    if zone_id is None:
        headers.append("منطقه")

    rows = []
    for m in members:
        row = [m["first_name"], m["last_name"], m["national_code"] or "", m["education"] or "",
               m["mobile"] or "", m["member_group"] or "", m["position"] or ""]
        if zone_id is None:
            row.append(zones_by_id.get(m["zone_id"], ""))
        rows.append(row)

    title = "گزارش اعضای شورای محلات"
    _write_excel_sheet(ws, title, headers, rows)
    wb.save(output_path)
    return output_path


def generate_requests_report_excel(db, output_path, zone_id=None):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "درخواست‌ها"

    requests = db.get_priority_requests(zone_id=zone_id)
    zones_by_id = {z["id"]: z["name"] for z in db.get_zones()}

    rows = []
    for i, r in enumerate(requests, start=1):
        rows.append([i, zones_by_id.get(r["zone_id"], ""), r["description"],
                     r["related_office"] or "", r["action_count"]])

    _write_excel_sheet(ws, "گزارش درخواست‌ها و مشکلات",
                        ["ردیف", "منطقه", "شرح درخواست/مشکل", "اداره مرتبط", "تعداد اقدام"], rows)
    wb.save(output_path)
    return output_path


def generate_actions_report_excel(db, output_path, zone_id=None):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "اقدامات"

    requests = db.get_priority_requests(zone_id=zone_id)
    zones_by_id = {z["id"]: z["name"] for z in db.get_zones()}

    rows = []
    for r in requests:
        actions = db.get_request_actions(r["id"])
        for a in actions:
            rows.append([
                zones_by_id.get(r["zone_id"], ""),
                r["description"],
                a["action_description"],
                a["created_at"]
            ])

    _write_excel_sheet(ws, "گزارش اقدامات انجام‌شده",
                        ["منطقه", "شرح درخواست/مشکل", "شرح اقدام", "تاریخ ثبت"], rows)
    wb.save(output_path)
    return output_path

# ==================== PowerPoint ====================

def _pptx_imports():
    """واردکردن وابستگی PowerPoint با پیام خطای قابل فهم."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        return Presentation, Inches, Pt, PP_ALIGN, RGBColor
    except Exception as exc:
        raise RuntimeError(
            "برای خروجی PowerPoint باید بسته python-pptx نصب باشد.\n"
            "در مسیر برنامه اجرا کنید: pip install python-pptx"
        ) from exc


def _ppt_font_name():
    # Tahoma روی بیشتر سیستم‌های ویندوز موجود است و برای فارسی مناسب‌تر از فونت‌های پیش‌فرض است.
    return "Tahoma"


def _ppt_add_textbox(slide, text, x, y, w, h, font_size=18, bold=False, color=(20, 40, 80), align_right=True):
    _, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT if align_right else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = str(text or "—")
    run.font.name = _ppt_font_name()
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def _ppt_add_title_slide(prs, title, subtitle=None):
    _, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(245, 248, 252)
    _ppt_add_textbox(slide, title, 0.6, 1.7, 12.0, 0.8, font_size=30, bold=True, color=(19, 41, 75))
    if subtitle:
        _ppt_add_textbox(slide, subtitle, 0.8, 2.7, 11.6, 0.5, font_size=16, color=(80, 92, 110))
    now_str = now_jalali()
    _ppt_add_textbox(slide, f"تاریخ تهیه گزارش: {now_str}", 0.8, 6.7, 11.6, 0.4, font_size=12, color=(100, 110, 125))
    return slide


def _ppt_add_section_header(slide, title, y=0.25):
    _ppt_add_textbox(slide, title, 0.5, y, 12.3, 0.45, font_size=20, bold=True, color=(19, 41, 75))


def _ppt_add_footer(slide, text="سامانه جامع مدیریت محلات و بلوک‌های شهرستان جوانرود"):
    _ppt_add_textbox(slide, text, 0.5, 7.1, 12.3, 0.25, font_size=9, color=(120, 130, 145))


def _ppt_add_rows_table(slide, headers, rows, x=0.6, y=1.0, w=12.1, row_h=0.34, max_rows=12):
    """جدول ساده RTL. ردیف‌های اضافی برش می‌خورند تا اسلاید شلوغ نشود."""
    _, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    display_rows = rows[:max_rows]
    cols = len(headers)
    table_shape = slide.shapes.add_table(len(display_rows) + 1, cols, Inches(x), Inches(y), Inches(w), Inches(row_h * (len(display_rows)+1))).table
    # عرض ستون‌ها مساوی؛ برای جداول دو ستونه ستون مقدار بزرگ‌تر است.
    for c in range(cols):
        table_shape.columns[c].width = Inches(w / cols)
    if cols == 2:
        table_shape.columns[0].width = Inches(w * 0.36)
        table_shape.columns[1].width = Inches(w * 0.64)
    for c, h in enumerate(headers):
        cell = table_shape.cell(0, c)
        cell.text = str(h)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(19, 41, 75)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.RIGHT
            for r in p.runs:
                r.font.name = _ppt_font_name(); r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = RGBColor(255,255,255)
    for r_idx, row in enumerate(display_rows, start=1):
        for c, val in enumerate(row):
            cell = table_shape.cell(r_idx, c)
            cell.text = str(val if val is not None else "—")
            if r_idx % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(240, 243, 248)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT
                for rr in p.runs:
                    rr.font.name = _ppt_font_name(); rr.font.size = Pt(9); rr.font.color.rgb = RGBColor(35, 45, 60)
    if len(rows) > max_rows:
        _ppt_add_textbox(slide, f"+ {len(rows) - max_rows} ردیف دیگر در فایل‌های PDF/Excel موجود است.", x, y + row_h*(max_rows+1) + 0.15, w, 0.28, font_size=9, color=(120, 80, 20))
    return table_shape


def _ppt_add_bullets(slide, items, x=0.8, y=1.1, w=11.5, h=5.6, font_size=14):
    _, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        p.level = 0
        run = p.add_run()
        run.text = "• " + str(item)
        run.font.name = _ppt_font_name()
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(35, 45, 60)
    return box


def _ppt_save_image_source(image_source):
    """اگر تصویر BytesIO باشد، آن را در فایل موقت ذخیره می‌کند و مسیر برمی‌گرداند."""
    if image_source is None:
        return None
    if isinstance(image_source, str):
        return image_source if os.path.exists(image_source) else None
    try:
        temp_dir = get_temp_dir()
        os.makedirs(temp_dir, exist_ok=True)
        path = os.path.join(temp_dir, "ppt_zone_snapshot.png")
        image_source.seek(0)
        with open(path, "wb") as f:
            f.write(image_source.read())
        return path
    except Exception:
        return None


def _ppt_add_zone_snapshot_slide(prs, db, zone_id, title="نمای گرافیکی بلوک", map_image_path=None):
    _, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, title)
    image_source, snapshot_meta = _snapshot_image_source(db, zone_id, map_image_path)
    image_path = _ppt_save_image_source(image_source)
    if image_path:
        # تصویر در مرکز و بزرگ نمایش داده می‌شود؛ برای پاورپوینت مهم‌ترین خروجی همین نقشه است.
        slide.shapes.add_picture(image_path, Inches(1.1), Inches(1.0), width=Inches(11.0), height=Inches(5.9))
        if snapshot_meta:
            _ppt_add_textbox(slide, f"نسخه تصویر: {snapshot_meta.get('version', 1)}", 1.1, 6.95, 11.0, 0.25, font_size=9, color=(95, 105, 120))
    else:
        _ppt_add_textbox(slide, "نمای گرافیکی بلوک هنوز تولید نشده است.", 1.0, 3.2, 11.0, 0.5, font_size=18, color=(160, 50, 50))
    _ppt_add_footer(slide)
    return slide


def generate_overview_report_pptx(db, output_path):
    Presentation, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _ppt_add_title_slide(prs, "گزارش کلی وضعیت سامانه", "سامانه جامع مدیریت محلات و بلوک‌های شهرستان جوانرود")

    zones = db.get_zones()
    all_streets = db.get_streets()
    all_places = db.get_places()
    all_members = db.get_council_members()
    all_requests = db.get_priority_requests()
    all_committees = [c for z in zones for c in get_zone_committee_report_data(db, z["id"])]
    total_actions = sum(r.get("action_count", 0) for r in all_requests)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "خلاصه آماری")
    rows = [
        ["مناطق/بلوک‌ها", len(zones)], ["خیابان‌ها و کوچه‌ها", len(all_streets)],
        ["اماکن", len(all_places)], ["مساجد مرجع", len(db.get_mosques())],
        ["اعضای شورا", len(all_members)], ["کمیته‌های تخصصی", len(all_committees)],
        ["اعضای فعال کمیته‌ها", sum(len(c.get("active_members") or []) for c in all_committees)],
        ["درخواست‌ها/مشکلات", len(all_requests)],
        ["مجموع اقدامات", total_actions],
    ]
    _ppt_add_rows_table(slide, ["شاخص", "مقدار"], rows, max_rows=10)
    _ppt_add_footer(slide)

    if zones:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_add_section_header(slide, "وضعیت بلوک‌ها")
        zone_rows = []
        for z in zones:
            zone_rows.append([
                z.get("name"), f"{(z.get('area_m2', 0) or 0)/10000:.2f}",
                len(db.get_streets(zone_id=z["id"])), len(db.get_mosques(zone_id=z["id"])),
                z.get("status") or "—",
            ])
        _ppt_add_rows_table(slide, ["بلوک", "مساحت هکتار", "معبر", "مسجد", "وضعیت"], zone_rows, max_rows=12)
        _ppt_add_footer(slide)

    prs.save(output_path)
    return output_path


def generate_block_full_report_pptx(db, zone_id, output_path, map_image_path=None):
    Presentation, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    zone = db.get_zone(zone_id)
    if not zone:
        raise ValueError("منطقه یافت نشد")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _ppt_add_title_slide(prs, f"گزارش کامل بلوک: {zone.get('name')}", "خروجی PowerPoint")
    _ppt_add_zone_snapshot_slide(prs, db, zone_id, title="۱. نمای گرافیکی بلوک", map_image_path=map_image_path)

    streets = db.get_streets(zone_id=zone_id)
    places = db.get_places(zone_id=zone_id)
    mosques = db.get_mosques(zone_id=zone_id)
    members = db.get_council_members(zone_id=zone_id)
    requests = db.get_priority_requests(zone_id=zone_id)
    meeting_place = db.get_zone_meeting_place(zone_id)
    profile = db.get_zone_profile(zone_id)
    neighborhood_issues = db.get_neighborhood_issues(zone_id)
    neighborhood_actions = db.get_neighborhood_actions(zone_id)
    neighborhood_meetings = db.get_neighborhood_meetings(zone_id)
    neighborhood_resolutions = db.get_neighborhood_resolutions(zone_id=zone_id)
    neighborhood_budgets = db.get_neighborhood_budgets(zone_id)
    budget_summary = db.get_budget_summary(zone_id)
    performance = db.get_zone_performance(zone_id)
    management_alerts = db.get_management_alerts(zone_id)
    field_visits = db.get_field_visits(zone_id)
    citizen_requests = db.get_citizen_requests(zone_id)
    operational_analysis = db.get_zone_operational_analysis(zone_id)
    correspondence_letters = db.get_correspondence_letters(zone_id=zone_id)
    approval_requests = db.get_approval_requests(zone_id=zone_id, limit=5000)
    generated_documents = db.get_generated_documents(zone_id=zone_id, limit=5000)
    committees = get_zone_committee_report_data(db, zone_id)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۲. مشخصات پایه بلوک")
    rows = [
        ["نام بلوک", zone.get("name") or "—"],
        ["وضعیت", zone.get("status") or "—"],
        ["مساحت", f"{(zone.get('area_m2', 0) or 0)/10000:.2f} هکتار"],
        ["محیط", f"{zone.get('perimeter_m', 0) or 0:.0f} متر"],
        ["تعداد معابر", len(streets)], ["تعداد اماکن", len(places)], ["تعداد مساجد", len(mosques)],
        ["کمیته‌های تخصصی", len(committees)], ["اعضای کمیته‌ها", sum(len(c.get("active_members") or []) for c in committees)],
        ["تعداد اعضا", len(members)], ["خانوار تأییدشده", profile.get("approved_households", 0)],
        ["جمعیت تخمینی", profile.get("estimated_population", 0)],
        ["مسائل باز", sum(1 for i in neighborhood_issues if i.get("status") not in ("مختومه", "انجام‌شده"))],
        ["بازدید میدانی", len(field_visits)], ["درخواست مردمی", len(citizen_requests)],
        ["مکاتبات اداری", len(correspondence_letters)],
        ["گردش تأیید", len(approval_requests)], ["اسناد Word", len(generated_documents)],
        ["ریسک عملیاتی", f"{operational_analysis.get('risk_score', 0)} — {operational_analysis.get('risk_level', '—')}"],
    ]
    _ppt_add_rows_table(slide, ["شاخص", "مقدار"], rows, max_rows=12)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۳. محل جلسه و مساجد")
    items = []
    if meeting_place:
        items.append(f"محل جلسه: {meeting_place.get('place_name') or '—'} — {meeting_place.get('exact_address') or '—'}")
    else:
        items.append("محل جلسه: ثبت نشده")
    if mosques:
        items.extend([f"مسجد: {m.get('name')}" for m in mosques[:10]])
    else:
        items.append("مسجدی برای این بلوک ثبت نشده است.")
    _ppt_add_bullets(slide, items, font_size=14)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۴. خیابان‌ها و کوچه‌ها")
    street_rows = [[s.get("name") or "بدون نام", s.get("highway_type") or "—"] for s in streets]
    _ppt_add_rows_table(slide, ["نام", "نوع"], street_rows or [["—", "—"]], max_rows=14)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۵. اماکن و اعضای شورا")
    place_items = [f"مکان: {p.get('name')} ({p.get('category') or p.get('type') or '—'})" for p in places[:6]]
    member_items = [f"عضو/معتمد: {m.get('full_name') or m.get('name') or '—'} — {m.get('phone') or '—'}" for m in members[:6]]
    _ppt_add_bullets(slide, (place_items + member_items) or ["موردی ثبت نشده است."], font_size=13)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۵-۱. کمیته‌های تخصصی بلوک")
    _ppt_add_rows_table(
        slide, ["ردیف", "کمیته", "رئیس", "دبیر", "اعضا", "جلسات", "مصوبات باز", "وضعیت"],
        committee_summary_rows(committees) or [["—", "کمیته‌ای ثبت نشده", "—", "—", "۰", "۰", "۰", "—"]],
        max_rows=8, row_h=0.46,
    )
    _ppt_add_footer(slide)

    for committee in committees:
        members_for_slide = committee.get("members") or []
        if not members_for_slide:
            members_for_slide = [None]
        for part_index in range(0, len(members_for_slide), 8):
            chunk = members_for_slide[part_index:part_index + 8]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            suffix = f" — بخش {part_index // 8 + 1}" if len(members_for_slide) > 8 else ""
            _ppt_add_section_header(slide, f"اعضا و سمت‌ها — {committee.get('title') or 'کمیته'}{suffix}")
            _ppt_add_bullets(slide, [
                f"رئیس: {committee.get('chair_name') or '—'}",
                f"دبیر: {committee.get('secretary_name') or '—'}",
                f"وضعیت: {committee.get('status') or '—'}",
                f"دستگاه‌های پیشنهادی: {committee.get('recommended_agencies') or '—'}",
            ], x=0.8, y=0.85, w=11.8, h=1.15, font_size=10)
            if chunk == [None]:
                rows = [["فاقد عضو ثبت‌شده", "—", "—", "—", "—"]]
            else:
                rows = [[
                    m.get("person_name") or "—", member_display_role(m), m.get("member_type") or "—",
                    m.get("agency_name") or "—", m.get("mobile") or "—",
                ] for m in chunk]
            _ppt_add_rows_table(
                slide, ["نام عضو", "سمت", "نوع عضویت", "اداره/دستگاه", "موبایل"],
                rows, x=0.65, y=2.1, w=12.0, row_h=0.48, max_rows=8,
            )
            _ppt_add_footer(slide)

    resolution_rows_all = committee_resolution_rows(committees)
    if resolution_rows_all:
        for part_index in range(0, len(resolution_rows_all), 9):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            suffix = f" — بخش {part_index // 9 + 1}" if len(resolution_rows_all) > 9 else ""
            _ppt_add_section_header(slide, f"مصوبات کمیته‌ها{suffix}")
            display_rows = [row[:6] for row in resolution_rows_all[part_index:part_index + 9]]
            _ppt_add_rows_table(
                slide, ["کمیته", "عنوان", "مسئول", "دستگاه", "مهلت", "وضعیت"],
                display_rows, max_rows=9, row_h=0.48,
            )
            _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۶. مسائل و نیازهای محله")
    issue_rows = [[i.get("title") or "—", i.get("priority_level") or "—", i.get("status") or "—", i.get("related_office") or "—"] for i in neighborhood_issues]
    _ppt_add_rows_table(slide, ["عنوان", "اولویت", "وضعیت", "دستگاه"], issue_rows or [["—", "—", "—", "—"]], max_rows=11)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۷. اقدامات اجرایی")
    action_rows = [[a.get("title") or "—", a.get("responsible_office") or a.get("responsible_person") or "—", a.get("status") or "—", f"{a.get('progress_percent') or 0}٪"] for a in neighborhood_actions]
    _ppt_add_rows_table(slide, ["عنوان", "مسئول", "وضعیت", "پیشرفت"], action_rows or [["—", "—", "—", "—"]], max_rows=11)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۸. جلسات و مصوبات")
    meeting_items = [f"جلسه: {m.get('title')} — {m.get('meeting_date') or '—'} — {m.get('status') or '—'}" for m in neighborhood_meetings[:6]]
    resolution_items = [f"مصوبه: {r.get('title')} — مسئول: {r.get('responsible_office') or r.get('responsible_person') or '—'} — {r.get('status') or '—'}" for r in neighborhood_resolutions[:7]]
    _ppt_add_bullets(slide, (meeting_items + resolution_items) or ["جلسه یا مصوبه‌ای ثبت نشده است."], font_size=12)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۹. مکاتبات اداری مرتبط")
    letter_rows = [[l.get("letter_number") or "—", l.get("direction") or "—", l.get("subject") or "—",
                    l.get("due_date") or "—", l.get("status") or "—"] for l in correspondence_letters]
    _ppt_add_rows_table(slide, ["شماره", "نوع", "موضوع", "مهلت", "وضعیت"],
                        letter_rows or [["—", "—", "موردی ثبت نشده", "—", "—"]], max_rows=10)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۱۰. گردش تأیید و اسناد اداری")
    approval_rows = [[a.get("title") or "—", a.get("entity_type") or "—",
                      f"{a.get('current_step') or 0}/{a.get('total_steps') or 0}",
                      a.get("due_date") or "—", a.get("status") or "—"] for a in approval_requests]
    _ppt_add_rows_table(slide, ["عنوان", "نوع", "مرحله", "مهلت", "وضعیت"],
                        approval_rows or [["—", "—", "—", "—", "موردی ثبت نشده"]], max_rows=8)
    if generated_documents:
        _ppt_add_bullets(slide, [f"سند: {d.get('title')} — {d.get('template_name') or 'قالب نامشخص'}" for d in generated_documents[:4]], font_size=10)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۱۱. عملیات میدانی")
    visit_rows = [[v.get("visit_date") or "—", v.get("officer_name") or "—", v.get("visit_type") or "—",
                   v.get("location_text") or "—", v.get("households_count") or 0, v.get("status") or "—"]
                  for v in field_visits]
    _ppt_add_rows_table(slide, ["تاریخ", "کارشناس", "نوع", "موقعیت", "خانوار", "وضعیت"],
                        visit_rows or [["—", "—", "—", "—", "۰", "—"]], max_rows=10)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۱۲. درخواست‌های مردمی و تحلیل عملیاتی")
    request_rows = [[r.get("tracking_code") or "—", r.get("title") or "—", r.get("urgency") or 0,
                     r.get("assigned_office") or "—", r.get("status") or "—"] for r in citizen_requests]
    _ppt_add_rows_table(slide, ["کد رهگیری", "عنوان", "فوریت", "دستگاه", "وضعیت"],
                        request_rows or [["—", "—", "۰", "—", "—"]], max_rows=8)
    _ppt_add_bullets(slide, [
        f"امتیاز ریسک: {operational_analysis.get('risk_score', 0)} — سطح: {operational_analysis.get('risk_level', '—')}",
        f"درخواست باز: {operational_analysis.get('open_requests', 0)} — درخواست فوری: {operational_analysis.get('urgent_requests', 0)}",
        f"مسئله بحرانی: {operational_analysis.get('critical_issues', 0)} — اقدام معوق: {operational_analysis.get('overdue_actions', 0)}",
    ], font_size=11)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۱۳. بودجه و هزینه")
    budget_rows = [[b.get("title") or "—", b.get("funding_source") or "—",
                    f"{b.get('allocated_amount') or 0:,.0f}", f"{b.get('spent_amount') or 0:,.0f}",
                    b.get("status") or "—"] for b in neighborhood_budgets]
    _ppt_add_rows_table(slide, ["عنوان", "منبع", "تخصیص", "هزینه", "وضعیت"],
                        budget_rows or [["—", "—", "۰", "۰", "—"]], max_rows=10)
    _ppt_add_footer(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "۱۴. ارزیابی عملکرد و هشدارها")
    items = [
        f"امتیاز کل: {performance.get('total_score', 0):.1f} از ۱۰۰ — سطح: {performance.get('level', '—')}",
        f"تکمیل اطلاعات: {performance.get('completeness', 0):.1f}٪",
        f"حل مسائل: {performance.get('issue_resolution', 0):.1f}٪",
        f"تکمیل اقدامات: {performance.get('action_completion', 0):.1f}٪",
        f"تحقق مصوبات: {performance.get('resolution_completion', 0):.1f}٪",
        f"پاسخ‌گویی و مشارکت: {performance.get('participation_response', 0):.1f}٪",
        f"اعتبار مصوب: {budget_summary.get('approved', 0):,.0f} ریال — هزینه: {budget_summary.get('spent', 0):,.0f} ریال",
    ]
    items.extend([f"هشدار {a.get('severity')}: {a.get('title')}" for a in management_alerts[:6]])
    _ppt_add_bullets(slide, items, font_size=13)
    _ppt_add_footer(slide)

    prs.save(output_path)
    return output_path


def generate_zone_full_report_pptx(db, zone_id, output_path):
    return generate_block_full_report_pptx(db, zone_id, output_path)


def generate_members_report_pptx(db, output_path, zone_id=None):
    Presentation, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    title = "گزارش اعضای شورای محلات" if zone_id is None else f"گزارش اعضای شورا — {db.get_zone(zone_id).get('name')}"
    _ppt_add_title_slide(prs, title, "خروجی PowerPoint")
    members = db.get_council_members(zone_id=zone_id)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "فهرست اعضا")
    rows = [[((m.get("first_name") or "") + " " + (m.get("last_name") or "")).strip() or "—", m.get("role") or "—", m.get("mobile") or "—"] for m in members]
    _ppt_add_rows_table(slide, ["نام", "نقش", "تلفن"], rows or [["—", "—", "—"]], max_rows=14)
    _ppt_add_footer(slide)
    prs.save(output_path); return output_path


def generate_requests_report_pptx(db, output_path, zone_id=None):
    Presentation, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    title = "گزارش درخواست‌ها و مشکلات" if zone_id is None else f"گزارش درخواست‌ها — {db.get_zone(zone_id).get('name')}"
    _ppt_add_title_slide(prs, title, "خروجی PowerPoint")
    requests = db.get_priority_requests(zone_id=zone_id)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "موارد اولویت‌دار")
    rows = []
    for r in requests:
        rows.append([r.get("title") or r.get("description") or "—", r.get("priority") or "—", r.get("status") or "—", r.get("action_count", 0)])
    _ppt_add_rows_table(slide, ["عنوان", "اولویت", "وضعیت", "اقدام"], rows or [["—", "—", "—", "—"]], max_rows=12)
    _ppt_add_footer(slide)
    prs.save(output_path); return output_path


def generate_actions_report_pptx(db, output_path, zone_id=None):
    Presentation, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    title = "گزارش اقدامات انجام‌شده" if zone_id is None else f"گزارش اقدامات — {db.get_zone(zone_id).get('name')}"
    _ppt_add_title_slide(prs, title, "خروجی PowerPoint")
    requests = db.get_priority_requests(zone_id=zone_id)
    actions = []
    for r in requests:
        for a in db.get_request_actions(r.get("id")):
            actions.append([r.get("title") or r.get("description") or "—", a.get("action_date") or "—", a.get("description") or a.get("title") or "—"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_section_header(slide, "اقدامات ثبت‌شده")
    _ppt_add_rows_table(slide, ["درخواست", "تاریخ", "شرح اقدام"], actions or [["—", "—", "—"]], max_rows=12)
    _ppt_add_footer(slide)
    prs.save(output_path); return output_path

# ==================== Correspondence Reports v6.4 ====================
def generate_correspondence_report_pdf(db, output_path, zone_id=None):
    title = "گزارش مکاتبات اداری"
    if zone_id is not None:
        zone = db.get_zone(zone_id)
        title += f" — {zone.get('name') if zone else zone_id}"
    letters = db.get_correspondence_letters(zone_id=zone_id, limit=10000)
    assignments = []
    for letter in letters:
        assignments.extend(db.get_workflow_assignments(letter_id=letter["id"], limit=1000))
    doc = SimpleDocTemplate(
        output_path, pagesize=A4, topMargin=12*mm, bottomMargin=15*mm,
        leftMargin=12*mm, rightMargin=12*mm,
    )
    flowables = _pdf_header_flowables(title)
    _, _, section_style, normal_style = _pdf_styles()
    flowables.append(Paragraph(shape_fa(f"دفتر مکاتبات (تعداد: {len(letters)})"), section_style))
    if letters:
        rows = [[
            l.get("letter_number") or "—", l.get("direction") or "—", l.get("subject") or "—",
            l.get("zone_name") or "—", l.get("sender") or "—", l.get("recipient") or "—",
            l.get("due_date") or "—", l.get("status") or "—", l.get("priority") or "—"
        ] for l in letters]
        flowables.append(_make_pdf_table(
            ["شماره", "نوع", "موضوع", "بلوک", "فرستنده", "گیرنده", "مهلت", "وضعیت", "اولویت"],
            rows, col_widths=[18*mm, 16*mm, 40*mm, 24*mm, 22*mm, 22*mm, 18*mm, 22*mm, 18*mm]
        ))
    else:
        flowables.append(Paragraph(shape_fa("مکاتبه‌ای ثبت نشده است."), normal_style))
    flowables.append(Spacer(1, 6*mm))
    flowables.append(Paragraph(shape_fa(f"ارجاعات و پیگیری‌ها (تعداد: {len(assignments)})"), section_style))
    if assignments:
        rows = [[
            a.get("letter_number") or "—", a.get("subject") or "—", a.get("assigned_to_name") or "—",
            a.get("instruction") or "—", a.get("due_date") or "—", a.get("priority") or "—",
            a.get("status") or "—", a.get("response_text") or "—"
        ] for a in assignments]
        flowables.append(_make_pdf_table(
            ["شماره نامه", "موضوع", "ارجاع به", "دستور", "مهلت", "اولویت", "وضعیت", "پاسخ"],
            rows, col_widths=[20*mm, 34*mm, 25*mm, 34*mm, 18*mm, 17*mm, 22*mm, 34*mm]
        ))
    else:
        flowables.append(Paragraph(shape_fa("ارجاعی ثبت نشده است."), normal_style))
    doc.build(flowables)
    return output_path


def generate_correspondence_report_excel(db, output_path, zone_id=None):
    letters = db.get_correspondence_letters(zone_id=zone_id, limit=10000)
    assignments = []
    for letter in letters:
        assignments.extend(db.get_workflow_assignments(letter_id=letter["id"], limit=1000))
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "مکاتبات"
    _write_excel_sheet(
        ws, "دفتر مکاتبات اداری",
        ["شناسه", "شماره", "نوع", "موضوع", "بلوک", "فرستنده", "گیرنده", "تاریخ نامه", "تاریخ دریافت", "مهلت", "وضعیت", "اولویت", "طبقه‌بندی", "پیوست", "ارجاع باز"],
        [[l.get("id"), l.get("letter_number"), l.get("direction"), l.get("subject"), l.get("zone_name"),
          l.get("sender"), l.get("recipient"), l.get("letter_date"), l.get("received_date"), l.get("due_date"),
          l.get("status"), l.get("priority"), l.get("confidentiality"), l.get("attachment_count"),
          l.get("open_assignment_count")] for l in letters],
    )
    ws2 = wb.create_sheet("ارجاعات")
    _write_excel_sheet(
        ws2, "کارتابل ارجاعات و پیگیری",
        ["شناسه", "شماره نامه", "موضوع", "بلوک", "ارجاع به", "دستور", "مهلت", "اولویت", "وضعیت", "پاسخ", "تاریخ تکمیل"],
        [[a.get("id"), a.get("letter_number"), a.get("subject"), a.get("zone_name"),
          a.get("assigned_to_name"), a.get("instruction"), a.get("due_date"), a.get("priority"),
          a.get("status"), a.get("response_text"), a.get("completed_at")] for a in assignments],
    )
    wb.save(output_path)
    return output_path


def generate_correspondence_report_pptx(db, output_path, zone_id=None):
    Presentation, Inches, Pt, PP_ALIGN, RGBColor = _pptx_imports()
    letters = db.get_correspondence_letters(zone_id=zone_id, limit=10000)
    assignments = []
    for letter in letters:
        assignments.extend(db.get_workflow_assignments(letter_id=letter["id"], limit=1000))
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    title = "گزارش مکاتبات اداری"
    if zone_id is not None:
        zone = db.get_zone(zone_id)
        title += f" — {zone.get('name') if zone else zone_id}"
    _ppt_add_title_slide(prs, title, f"{len(letters)} نامه و {len(assignments)} ارجاع")
    for start in range(0, max(1, len(letters)), 10):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_add_section_header(slide, "دفتر مکاتبات اداری")
        chunk = letters[start:start+10]
        rows = [[l.get("letter_number") or "—", l.get("direction") or "—", l.get("subject") or "—",
                 l.get("zone_name") or "—", l.get("due_date") or "—", l.get("status") or "—"] for l in chunk]
        _ppt_add_rows_table(slide, ["شماره", "نوع", "موضوع", "بلوک", "مهلت", "وضعیت"],
                            rows or [["—", "—", "موردی ثبت نشده", "—", "—", "—"]], max_rows=10)
        _ppt_add_footer(slide)
    for start in range(0, max(1, len(assignments)), 10):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_add_section_header(slide, "کارتابل ارجاعات")
        chunk = assignments[start:start+10]
        rows = [[a.get("letter_number") or "—", a.get("assigned_to_name") or "—",
                 a.get("instruction") or "—", a.get("due_date") or "—", a.get("status") or "—"] for a in chunk]
        _ppt_add_rows_table(slide, ["نامه", "ارجاع به", "دستور", "مهلت", "وضعیت"],
                            rows or [["—", "—", "موردی ثبت نشده", "—", "—"]], max_rows=10)
        _ppt_add_footer(slide)
    prs.save(output_path)
    return output_path
